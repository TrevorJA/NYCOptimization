import pathlib
# -*- coding: utf-8 -*-
"""Build docs/NYCOptimization_project_proposal_and_motivation_CLAUDE.pptx.

Content is the FINALIZED method: three scenario designs (historic,
fixed_probabilistic, hazard_filling) in one stationary population (Kirsch-Nowak
fit to the record, no climate perturbation); deep uncertainty enters ONLY in
the held-out test ensemble E_test. A deterministic LHS + nearest-neighbour
selector in ABSOLUTE, range-scaled hazard space (there is no simulated
annealing and no empirical-CDF/rank scaling); one budget condition at N = 100,
L = 10 yr; and the held-out metric set (multivariate satisficing primary, no
regret and no perfect-foresight optimization anywhere). The single controlled
contrast is fixed_probabilistic -> hazard_filling: same generator, population,
N, L, NFE, only the selection rule differs.

Sources of truth: docs/manuscript/Amestoy_NYC_reoptimization_manuscript_draft.md,
docs/notes/methods/{scenario_design_methods, experimental_design,
objective_definitions}.md, docs/notes/terminology.md, src/scenario_designs.py.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

INK    = RGBColor(0x21, 0x25, 0x29)
ACCENT = RGBColor(0x2E, 0x5E, 0x8C)
GRAY   = RGBColor(0x5F, 0x6B, 0x7A)
LIGHTB = RGBColor(0xE9, 0xEF, 0xF5)
RUST   = RGBColor(0xA4, 0x59, 0x29)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MIDGRAY= RGBColor(0x9A, 0xA4, 0xAF)

TITLE_FONT = "Cambria"
BODY_FONT  = "Calibri"
# Slide images, versioned in-repo so the deck is reproducible. These previously
# pointed at a session scratchpad under %TEMP%, which meant the authoritative
# deck silently stopped building whenever that directory was cleaned.
MEDIA = str(pathlib.Path(__file__).resolve().parents[2] / "docs" / "assets" / "proposal_media")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide(notes=None):
    s = prs.slides.add_slide(BLANK)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _apply_runs(p, text, size, color, font=BODY_FONT, italic=False, base_bold=False):
    parts = text.split("**")
    for i, part in enumerate(parts):
        if not part:
            continue
        r = p.add_run()
        r.text = part
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.italic = italic
        r.font.bold = base_bold or (i % 2 == 1)


def add_title(s, text, size=28):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12.13), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    _apply_runs(tf.paragraphs[0], text, size, INK, font=TITLE_FONT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.63), Inches(1.18), Inches(1.9), Pt(2.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    return tb


def add_body(s, items, left=0.6, top=1.5, width=12.13, height=None,
             size=20, space_after=10, line_spacing=1.04):
    if height is None:
        height = 7.25 - top
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for kind, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if kind == "gap":
            p.space_after = Pt(text if isinstance(text, (int, float)) else 6)
            continue
        if kind == "b":
            _apply_runs(p, u"•  " + text, size, INK)
        elif kind == "p":
            _apply_runs(p, text, size, INK)
        elif kind == "h":
            _apply_runs(p, text, size, INK, base_bold=True)
        elif kind == "c":
            _apply_runs(p, text, 18, GRAY, italic=True)
        elif kind == "f":
            _apply_runs(p, text, size, INK, font=TITLE_FONT)
    return tb


def box(s, x, y, w, h, text=None, fill=LIGHTB, line=GRAY, size=18, bold=False,
        color=INK, shape=MSO_SHAPE.RECTANGLE, align=PP_ALIGN.CENTER, line_w=1.0,
        font=BODY_FONT):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if text is not None:
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        _apply_runs(p, text, size, color, font=font, base_bold=bold)
    return sh


def label(s, x, y, w, h, text, size=18, color=GRAY, align=PP_ALIGN.CENTER,
          italic=False, bold=False, font=BODY_FONT, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    _apply_runs(p, text, size, color, font=font, italic=italic, base_bold=bold)
    return tb


def arrow(s, x1, y1, x2, y2, color=GRAY, weight=1.5):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def dot(s, x, y, d=0.09, color=ACCENT):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def ring(s, x, y, d=0.13, color=RUST, weight=1.25):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.background()
    sh.line.color.rgb = color; sh.line.width = Pt(weight)
    sh.shadow.inherit = False
    return sh


def line(s, x1, y1, x2, y2, color=MIDGRAY, weight=0.75, dash=None):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    conn.shadow.inherit = False
    if dash:
        ln = conn.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    return conn


def axes_panel(s, x, y, w, h, xlabel, ylabel):
    arrow(s, x, y + h, x + w, y + h, color=MIDGRAY, weight=1.2)
    arrow(s, x, y + h, x, y, color=MIDGRAY, weight=1.2)
    if xlabel:
        label(s, x, y + h + 0.03, w, 0.32, xlabel, size=18, color=GRAY)
    if ylabel:
        tb = label(s, x - 1.15, y + h / 2 - 0.175, 1.6, 0.35, ylabel, size=18, color=GRAY)
        tb.rotation = 270


def pic(s, name, x, y, w=None, h=None):
    kw = {}
    if w is not None: kw['width'] = Inches(w)
    if h is not None: kw['height'] = Inches(h)
    return s.shapes.add_picture(MEDIA + "/" + name, Inches(x), Inches(y), **kw)


NO_STYLE_TABLE = "{5940675A-B579-460E-94D1-54222C63F5DA}"

def add_table(s, rows, col_widths, left=0.6, top=1.5, row_h=0.5,
              size=18, header_fill=LIGHTB, highlight_row=None, highlight_color=None):
    n_rows, n_cols = len(rows), len(rows[0])
    gfx = s.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                             Inches(sum(col_widths)), Inches(row_h * n_rows))
    tbl = gfx.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('bandRow', '0'); tblPr.set('firstRow', '1')
    for el in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(el)
    style_el = tblPr.makeelement(qn('a:tableStyleId'), {})
    style_el.text = NO_STYLE_TABLE
    tblPr.append(style_el)
    for j, wdt in enumerate(col_widths):
        tbl.columns[j].width = Inches(wdt)
    highlight = highlight_row if isinstance(highlight_row, (list, tuple)) else (
        [] if highlight_row is None else [highlight_row])
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(row_h)
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            elif i in highlight:
                cell.fill.solid(); cell.fill.fore_color.rgb = (highlight_color or LIGHTB)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            _apply_runs(tf.paragraphs[0], cell_text, size, INK, base_bold=(i == 0))
    return gfx


def divider(s, text_lines):
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = LIGHTB
    band.line.fill.background(); band.shadow.inherit = False
    tf = band.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, bold) in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _apply_runs(p, txt, 32, INK, font=TITLE_FONT, base_bold=bold)


def claim(s, text, top=1.42):
    return label(s, 0.6, top, 12.13, 0.62, text, size=20, color=INK,
                 align=PP_ALIGN.LEFT, bold=True)


import random as _r

# =====================================================================
# 1 — Title
# =====================================================================
s = add_slide(notes=(
    "Working title; alternates: 'Designing the Scenarios We Optimize Over'. Audience: Pat Reed / "
    "committee. No results yet; deliverable is scope, gap, and planned experiment."))
label(s, 0.9, 2.15, 11.5, 1.6,
      "Hazard-Informed Scenario Design for\nMany-Objective Reservoir Policy Search",
      size=36, color=INK, align=PP_ALIGN.LEFT, font=TITLE_FONT)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(3.62), Inches(2.6), Pt(3))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
label(s, 0.9, 3.85, 11.5, 0.5,
      "Re-optimizing NYC reservoir operations in the Delaware River Basin under hydrologic uncertainty",
      size=22, color=GRAY, align=PP_ALIGN.LEFT)
label(s, 0.9, 5.35, 11.5, 1.2,
      "Research proposal  |  Trevor Amestoy\nReed Research Group, Cornell University  |  Summer 2026",
      size=20, color=INK, align=PP_ALIGN.LEFT)

# =====================================================================
# 2 — Contents
# =====================================================================
s = add_slide(notes="Roadmap; supplemental holds reference detail (objectives, forcing, scaling).")
add_title(s, "Contents")
add_body(s, [
    ("h", "Motivation & literature framing"),
    ("h", "The gap"),
    ("h", "Proposal"),
    ("b", "hypothesis  ·  research questions  ·  contributions"),
    ("h", "Planned experiment"),
    ("b", "three designs, one stationary population  ·  deep uncertainty only in re-evaluation  ·  the controlled contrast  ·  objectives  ·  controls"),
    ("h", "How designs are compared"),
    ("b", "held-out metric set  ·  threshold sweep  ·  mechanism test  ·  threats to validity"),
    ("h", "Status & open decisions"),
], left=0.9, top=1.7, width=5.7, size=22)
ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.9), Inches(1.7), Inches(6.9), Inches(6.3))
ln.line.color.rgb = MIDGRAY; ln.line.width = Pt(1.0); ln.shadow.inherit = False
add_body(s, [
    ("h", "Supplemental"),
    ("b", "objective formulations (7) and notation"),
    ("b", "forcing parameterization (CMIP6 harmonic)"),
    ("b", "Anvil scaling experiments"),
    ("b", "planned manuscript figure sequence"),
], left=7.4, top=1.7, width=5.3, size=22)

# =====================================================================
# 3 — Motivation 1: beyond the record
# =====================================================================
s = add_slide(notes=(
    "Left figure (from Jazmine's dissertation): FDC envelope keeps widening with ensemble size. "
    "Right (our DRB Kirsch-Nowak diagnostic): max drought magnitude and hazard hull area still "
    "growing past 10^4 scenario-years. TODO before presenting: annotate the dashed reference line "
    "and un-truncate the y-axis label in the right panels."))
add_title(s, "Motivation: the observed record undersamples hazard")
claim(s, "New hydrologic extremes keep emerging as stochastic ensembles grow.")
add_body(s, [
    ("b", "optimizing over one record overfits a single hydrologic history (Brodeur et al. 2020)"),
    ("b", "synthetic ensembles are the standard remedy (Kirsch et al. 2013; Herman et al. 2016)"),
], top=2.05, size=20, space_after=6)
pic(s, "image1.png", 0.9, 3.25, w=3.3)
label(s, 0.75, 6.85, 3.6, 0.4, "FDC range vs ensemble size", size=18, color=GRAY)
pic(s, "image2.png", 5.15, 3.6, w=7.5)
label(s, 5.15, 6.85, 7.5, 0.4,
      "drought extremes and hazard diversity vs scenario-years (this study)", size=18, color=GRAY)

# =====================================================================
# 4 — Motivation 2: low-signal years
# =====================================================================
s = add_slide(notes=(
    "Our current DRB stochastic ensembles: 84-93% of years never leave the Normal/Flood storage "
    "zone; simulating them costs the same as stress years but adds little gradient for the supply/"
    "drought objectives. If asked: the three legend percentages correspond to the three current "
    "stochastic ensemble presets."))
add_title(s, "Motivation: most simulated years carry little search signal")
claim(s, "84–93% of stochastic-ensemble years sit in Normal / Flood conditions.")
add_body(s, [
    ("b", "ensemble size trades off directly against evaluation cost in search (Zatarain Salazar et al. 2017)"),
    ("b", "supply and drought objectives draw signal only from the rare stress years"),
    ("b", "high-flow years are sparse in a complementary way; both motivate designed coverage"),
], left=0.6, top=2.5, width=6.4, size=20)
pic(s, "image3.png", 7.3, 1.95, w=5.4)
label(s, 7.15, 6.35, 5.7, 0.75, "share of years whose lowest NYC storage stays in the Normal / Flood FFMP zones, three ensembles (this study)", size=18, color=GRAY)

# =====================================================================
# 5 — Motivation 3: input-space redundancy
# =====================================================================
s = add_slide(notes=(
    "MOTIVATION ONLY: why hazard space is the interesting space to control, not a campaign contrast. "
    "Sampling generator/climate parameters (the prevailing deep-uncertainty design) does not control "
    "where realizations land in hazard space: the parameter-to-stress mapping is many-to-one. Right "
    "figure: a large forcing modification barely moves the flow distribution. Preliminary DRB "
    "diagnostic: input-parameter sampling leaves ~1.3x higher discrepancy (unevenness) in hazard "
    "space, surviving input-sample enrichment. This motivates selecting directly in hazard space; "
    "input-parameter sampling is NOT one of the three campaign designs."))
add_title(s, "Motivation: input-parameter sampling leaves hazard-space gaps")
claim(s, "Distinct generator parameters produce overlapping hazard conditions.")
add_body(s, [
    ("b", "sampling generator / climate parameters is the prevailing deep-uncertainty design (Quinn et al. 2020; Bartholomew & Kwakkel 2020)"),
    ("b", "distinct parameter sets often yield hydrologically redundant realizations (Guo et al. 2018; Quinn et al. 2020)"),
    ("b", "preliminary DRB diagnostic: sampling input parameters leaves about 1.3× higher discrepancy (unevenness) in hazard space, surviving input enrichment (this study)"),
    ("b", "this motivates controlling coverage **directly in hazard space** — the point of the proposed selector, not a campaign design of its own"),
], top=2.1, size=20)
py = 4.85; ph = 1.55; pw = 2.5
axes_panel(s, 2.5, py, pw, ph, "forcing parameter 1", "parameter 2")
for (fx, fy) in [(0.15, 0.2), (0.5, 0.75), (0.8, 0.35), (0.3, 0.55), (0.65, 0.12), (0.9, 0.8), (0.12, 0.85)]:
    dot(s, 2.5 + fx * (pw - 0.25) + 0.06, py + fy * (ph - 0.25) + 0.06, d=0.11, color=ACCENT)
arrow(s, 5.3, py + ph / 2, 6.1, py + ph / 2, color=GRAY, weight=1.8)
axes_panel(s, 6.35, py, pw, ph, "hazard metric 1", None)
for (ox, oy, ow, oh) in [(0.5, 0.45, 1.25, 0.7), (0.7, 0.3, 1.25, 0.7), (0.4, 0.6, 1.25, 0.7)]:
    e = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.35 + ox), Inches(py + oy - 0.15), Inches(ow), Inches(oh))
    e.fill.solid(); e.fill.fore_color.rgb = LIGHTB
    e.line.color.rgb = ACCENT; e.line.width = Pt(1.0); e.shadow.inherit = False
pic(s, "image5.png", 9.7, 4.35, w=3.3)
label(s, 9.7, 6.95, 3.3, 0.4, "modified forcing, similar flows", size=18, color=GRAY)

# =====================================================================
# 6 — Motivation 4: rare corners are consequential
# =====================================================================
s = add_slide(notes=(
    "Figure (ours): fraction of droughts, by severity x magnitude bin, that push NYC into the "
    "Drought Emergency zone; consequential outcomes concentrate in rare severe corners. Cohen 2021: "
    "high-baseline-regret (stressful) training subsets matched all-scenario training at ~3x less "
    "compute -- cited as MOTIVATION only; we never run a perfect-foresight optimization. Zaniolo "
    "2023: policies trained per drought type differ in structure and performance."))
add_title(s, "Motivation: rare hazard corners drive consequential outcomes")
claim(s, "Consequential failures concentrate in rare, severe drought corners.")
add_body(s, [
    ("b", "policies only learn from scenarios that contain these events"),
    ("b", "high-stress training scenarios matched full-ensemble robustness at about 3× less compute (Cohen et al. 2021)"),
    ("b", "training on distinct drought types changes the policies found (Zaniolo et al. 2023)"),
    ("b", "Cohen selects scenarios by a **perfect-foresight optimization per scenario**; we cite the finding, never the machinery"),
], left=0.6, top=2.5, width=7.4, size=20)
pic(s, "image8.png", 8.6, 1.85, h=4.95)
label(s, 8.2, 6.85, 5.0, 0.4, "this study: DRB stationary stochastic baseline", size=18, color=GRAY)

# =====================================================================
# 7 — Literature framing: four families (taxonomy w/ mini graphics)
# =====================================================================
s = add_slide(notes=(
    "Families defined by the criterion of a good scenario set: expert judgment (I), fidelity to a "
    "distribution (II), coverage of a condition space (III), relevance to the decision problem (IV). "
    "This study sits in family III, coverage-based subset selection, but applies the coverage "
    "criterion in hazard space measured on realized sequences. Full taxonomy in the working notes."))
add_title(s, "How the literature designs scenario sets: four families")
label(s, 0.6, 1.4, 12.13, 0.45,
      "Families defined by the criterion of a good scenario set (working taxonomy, after Fairbrother et al. 2022; Brown et al. 2012)",
      size=18, color=GRAY, align=PP_ALIGN.LEFT)
panel_y, panel_h, panel_w = 2.35, 1.45, 2.35
xs = [0.85, 4.0, 7.15, 10.3]
# I. judgment: a single trace
axes_panel(s, xs[0], panel_y, panel_w, panel_h, None, None)
trace = [(0.02, 0.55), (0.18, 0.2), (0.32, 0.65), (0.5, 0.1), (0.63, 0.5), (0.8, 0.3), (0.97, 0.6)]
for k in range(len(trace) - 1):
    x1, y1 = trace[k]; x2, y2 = trace[k + 1]
    line(s, xs[0] + x1 * panel_w, panel_y + y1 * panel_h, xs[0] + x2 * panel_w, panel_y + y2 * panel_h,
         color=ACCENT, weight=1.75)
# II. distribution-driven: central cluster
axes_panel(s, xs[1], panel_y, panel_w, panel_h, None, None)
_r.seed(3)
for _ in range(26):
    fx = 0.5 + 0.9 * (_r.random() - 0.5) * _r.random()
    fy = 0.5 + 0.9 * (_r.random() - 0.5) * _r.random()
    dot(s, xs[1] + 0.08 + fx * (panel_w - 0.3), panel_y + 0.05 + fy * (panel_h - 0.3), d=0.08, color=ACCENT)
# III. coverage: uniform grid
axes_panel(s, xs[2], panel_y, panel_w, panel_h, None, None)
for gx in range(5):
    for gy in range(3):
        dot(s, xs[2] + 0.22 + gx * (panel_w - 0.55) / 4, panel_y + 0.14 + gy * (panel_h - 0.5) / 2,
            d=0.09, color=ACCENT)
# IV. problem-driven: scattered, a few selected
axes_panel(s, xs[3], panel_y, panel_w, panel_h, None, None)
pts = [(0.15, 0.25), (0.35, 0.6), (0.5, 0.2), (0.62, 0.75), (0.8, 0.4), (0.25, 0.8),
       (0.7, 0.1), (0.9, 0.65), (0.45, 0.45), (0.85, 0.9)]
for i, (fx, fy) in enumerate(pts):
    sel = i in (3, 6, 9)
    dot(s, xs[3] + 0.08 + fx * (panel_w - 0.3), panel_y + 0.05 + fy * (panel_h - 0.3),
        d=0.12 if sel else 0.08, color=RUST if sel else MIDGRAY)
caps = [
    "**I. Judgment-based**\nsingle traces, design events\n(Giuliani et al. 2016)",
    "**II. Distribution-driven**\nsample an assumed distribution\n(Quinn et al. 2017)",
    "**III. Bottom-up coverage**\nspan a condition space\n(Brown et al. 2012; Quinn et al. 2020)",
    "**IV. Problem-driven**\nselect by decision relevance\n(Cohen et al. 2021)",
]
for x0, cap in zip(xs, caps):
    tb = s.shapes.add_textbox(Inches(x0 - 0.25), Inches(panel_y + panel_h + 0.12), Inches(panel_w + 0.55), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln_txt in enumerate(cap.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _apply_runs(p, ln_txt, 18, INK if i == 0 else GRAY)
add_body(s, [
    ("h", "This study: family III coverage, applied in hazard space measured on realized streamflow sequences."),
], top=5.75, size=20)

# =====================================================================
# 8 — Closest antecedents & the gap
# =====================================================================
s = add_slide(notes=(
    "The honest gap statement. Cohen 2021 already showed composition matters (cited as motivation, "
    "never as a metric we compute); each antecedent stops short on one axis: problem-driven vs "
    "simulation-free selection (Cohen), post-hoc vs search-phase (Bonham), generation control vs "
    "selection from realized sequences (Zaniolo), 1-D vs multi-dimensional (Zatarain Salazar). "
    "'Genuinely held-out' contrasts with Cohen's complementary-halves testing."))
add_title(s, "Closest antecedents and the gap")
rows = [
    ["Antecedent", "Established", "Stops short at"],
    ["Cohen et al. 2021", "training-scenario properties drive out-of-sample robustness", "problem-driven regret selection; a perfect-foresight run per scenario; 97 GCM traces"],
    ["Bonham et al. 2024", "space-filling subsampling + coverage metrics in water resources", "post-hoc robustness ranking, not the search ensemble"],
    ["Zaniolo et al. 2023", "drought-type training ensembles change the policies found", "imposed at generation; four discrete types"],
    ["Zatarain Salazar et al. 2017", "search-ensemble size vs fidelity vs cost tradeoff", "one-dimensional flow-magnitude stratification"],
]
add_table(s, rows, [2.75, 4.85, 4.55], left=0.6, top=1.5, row_h=0.68, size=18)
gb = box(s, 0.6, 5.2, 12.15, 1.35, None, fill=LIGHTB, line=ACCENT, line_w=1.25)
tf = gb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_apply_runs(tf.paragraphs[0],
            "**Gap.**  No study builds the MOEA search ensemble by space-filling selection in a multi-dimensional hazard space, at candidate-pool scale (10⁵–10⁶ sequences), and tests it on a held-out ensemble spanning deep climate uncertainty.",
            20, INK)
label(s, 0.6, 6.65, 12.15, 0.55,
      "Cohen et al. (2021) showed composition matters; the new element is the scalable, simulation-free construction and its held-out test.",
      size=18, color=GRAY, align=PP_ALIGN.LEFT, italic=True)

# =====================================================================
# 9 — Three spaces
# =====================================================================
s = add_slide(notes=(
    "Load-bearing vocabulary. Hazard space: drought event metrics (SSI run theory), low/high-flow "
    "indices, measured on each realized sequence before any simulation. Because coordinates exist "
    "pre-simulation, hazard-space ensemble design is simulation-free and reusable across "
    "formulations; that is what separates it from family IV problem-driven selection, which needs "
    "an optimization per scenario."))
add_title(s, "Terminology: three spaces")
dy = 2.3
box(s, 0.7, dy, 3.5, 1.1, "**Input space**\ngenerator & climate-forcing parameters θ", size=18)
arrow(s, 4.2, dy + 0.55, 4.95, dy + 0.55)
label(s, 3.68, dy + 1.18, 1.8, 0.35, "generate", size=18, color=GRAY)
box(s, 4.95, dy, 3.5, 1.1, "**Hazard space**\ndrought / low-flow / high-flow metrics per sequence", size=18, line=ACCENT, line_w=1.75)
arrow(s, 8.45, dy + 0.55, 9.2, dy + 0.55)
label(s, 7.93, dy + 1.18, 1.8, 0.35, "simulate", size=18, color=GRAY)
box(s, 9.2, dy, 3.5, 1.1, "**Outcome space**\nobjective values of a policy under a scenario", size=18)
add_body(s, [
    ("b", "hazard coordinates are **measured on realized sequences, before any simulation** (drought event metrics per run theory, Yevjevich 1967; SSI, Vicente-Serrano et al. 2012; flow indices, Richter et al. 1996)"),
    ("b", "so hazard-space ensemble design is a **simulation-free, pre-optimization step**, reusable across problem formulations"),
], top=4.3, size=20)

# =====================================================================
# 10 — The proposed design + the intrinsic asymmetry
# =====================================================================
s = add_slide(notes=(
    "THE central conceptual slide. You cannot generate a realization with prescribed hazard "
    "coordinates -- hazard is an emergent property of a realized flow sequence, not a knob on the "
    "generator. Forcing parameters theta ARE a knob. So input-space designs GENERATE TO their "
    "design points (LHS alone, nothing to snap to) while hazard-space designs must SELECT FROM a "
    "finite candidate pool (LHS anchors + nearest-neighbour snap). That asymmetry is intrinsic, not "
    "an implementation shortcut, and it is why hazard-filling is the only design that needs a pool. "
    "The hazard-filling design owns its i.i.d. candidate pool; nothing is shared. The selector fills "
    "ABSOLUTE, range-scaled hazard space, deliberately over-representing the severe corners."))
add_title(s, "Proposed design: hazard-filling search ensembles")
claim(s, "θ is a knob on the generator. Hazard is not — it is emergent, measured after the fact.")
add_body(s, [
    ("b", "input-parameter designs **generate to** their design points: LHS alone, nothing to snap to"),
    ("b", "hazard-space design must **select from** a finite **candidate pool** of 10⁵–10⁶ short (10-yr) i.i.d. realizations, which the design generates and owns"),
    ("b", "selector: **LHS anchors in ABSOLUTE, range-scaled hazard space, snapped to the nearest unused pool member** — deliberately over-represents the severe corners; deterministic, no tuning"),
    ("b", "the nearest-neighbour step is **intrinsic to hazard-space design**, not an approximation of something better; the asymmetry is part of the argument"),
    ("b", "the shift toward the severe corners is the deliberate intervention RQ1 tests; the single comparison point is the held-out, deeply uncertain re-evaluation on E_test"),
], top=2.05, size=20, space_after=7)
py = 5.15; ph = 1.75; pw = 2.85
axes_panel(s, 2.6, py, pw, ph, "hazard metric 1", "metric 2")
_r.seed(21)
for _ in range(70):
    fx, fy = _r.random(), _r.random()
    dot(s, 2.6 + 0.1 + fx * (pw - 0.35), py + 0.08 + fy * (ph - 0.35), d=0.055, color=MIDGRAY)
label(s, 2.1, py - 0.45, 3.9, 0.36, "candidate pool (own, i.i.d., ~10⁶)", size=18, color=GRAY)
arrow(s, 5.95, py + ph / 2, 7.15, py + ph / 2, color=GRAY, weight=1.8)
label(s, 5.35, py + ph / 2 - 0.92, 2.7, 0.7, "LHS anchors +\nnearest-neighbour snap", size=18, color=GRAY)
axes_panel(s, 9.2, py, pw, ph, "hazard metric 1", "metric 2")
for fx, fy in [(0.1, 0.15), (0.5, 0.1), (0.9, 0.15), (0.1, 0.5), (0.5, 0.5), (0.88, 0.52),
               (0.12, 0.85), (0.5, 0.88), (0.9, 0.85), (0.3, 0.3), (0.7, 0.32), (0.3, 0.7), (0.7, 0.7)]:
    dot(s, 9.2 + 0.1 + fx * (pw - 0.35), py + 0.08 + fy * (ph - 0.35), d=0.1, color=ACCENT)
label(s, 8.85, py - 0.45, 3.6, 0.36, "search ensemble (N = 100)", size=18, color=GRAY)

# =====================================================================
# 11 — Hypothesis & contributions
# =====================================================================
s = add_slide(notes=(
    "Falsifiable hypothesis; contribution 1 stands even under a null result. Positioning vs Cohen "
    "2021: they showed composition matters; we contribute the scalable, simulation-free, "
    "coverage-designed construction, an exact statistical control (fixed_probabilistic), and a "
    "genuinely held-out deep-uncertainty test."))
add_title(s, "Hypothesis & planned contributions")
hb = box(s, 0.6, 1.6, 12.15, 1.45, None, fill=LIGHTB, line=ACCENT, line_w=1.25)
tf = hb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_apply_runs(tf.paragraphs[0],
            "**Hypothesis.**  Search ensembles that fill hazard space produce Pareto-approximate policies that are more robust under held-out re-evaluation than probabilistic (i.i.d.) sampling from the same stochastic generator, at equal N, equal L and equal NFE.",
            20, INK)
add_body(s, [
    ("h", "Planned contributions"),
    ("b", "**first budget-controlled comparison of scenario designs**, with an **exact statistical control** (fixed_probabilistic) for the hazard-filling claim, judged on a genuinely held-out, deeply uncertain re-evaluation"),
    ("b", "**scalable, simulation-free, coverage-designed** search-ensemble construction (hazard-filling), selecting in absolute hazard space over its own i.i.d. pool"),
    ("b", "**design-ranking stability under threshold stringency** — the sweep nobody in this lineage has run"),
    ("b", "evidence on **FFMP re-optimization** and a variable-resolution rule structure"),
], top=3.35, size=20, space_after=8)

# =====================================================================
# 12 — Testbed: NYC reservoirs in the DRB
# =====================================================================
s = add_slide(notes=(
    "One context slide so the objectives and RQ2/RQ3 are interpretable to non-DRB readers: three "
    "NYC Delaware reservoirs, the 1954 Supreme Court Decree quantities (diversion cap, Montague "
    "and Trenton flow targets), and the FFMP acronym expanded. All simulation in Pywr-DRB."))
add_title(s, "Testbed: NYC reservoirs in the Delaware River Basin")
add_body(s, [
    ("b", "Cannonsville, Pepacton and Neversink supply about half of NYC's water; the 1954 Supreme Court Decree caps the diversion at 800 MGD"),
    ("b", "the Decree also sets minimum-flow targets at Montague and Trenton, protecting New Jersey and Philadelphia (salinity repulsion)"),
    ("b", "operations follow the **Flexible Flow Management Program (FFMP)**: storage-zone-based releases with drought step-downs; its 24 rule parameters are the decision variables re-optimized here"),
    ("b", "simulated with Pywr-DRB, an open-source daily water-availability model of the basin (Hamilton et al. 2024)"),
], left=0.6, top=1.55, width=6.3, size=20)
rx = 8.55
box(s, rx + 0.0, 1.7, 1.5, 0.62, "Cannonsville", size=18)
box(s, rx + 1.62, 1.7, 1.5, 0.62, "Pepacton", size=18)
box(s, rx + 3.24, 1.7, 1.5, 0.62, "Neversink", size=18)
arrow(s, rx + 0.75, 2.32, rx + 2.2, 3.15)
arrow(s, rx + 2.37, 2.32, rx + 2.37, 3.15)
arrow(s, rx + 3.99, 2.32, rx + 2.55, 3.15)
box(s, rx + 1.45, 3.2, 1.9, 0.6, "Delaware mainstem", fill=None, line=MIDGRAY, size=18)
arrow(s, rx + 2.4, 3.8, rx + 2.4, 4.35)
box(s, rx + 1.2, 4.4, 2.3, 0.72, "Montague target\n1,750 cfs", size=18)
arrow(s, rx + 2.35, 5.12, rx + 2.35, 5.65)
box(s, rx + 1.2, 5.7, 2.3, 0.72, "Trenton target\n3,000 cfs", size=18)
box(s, 7.05, 2.75, 1.6, 0.85, "NYC diversion\n≤ 800 MGD", fill=WHITE, line=ACCENT, size=18)
arrow(s, rx + 0.5, 2.35, 8.3, 2.8)

# =====================================================================
# 13 — Research questions
# =====================================================================
s = add_slide(notes="RQ1 drives the experimental design; RQ2/RQ3 answered by the same campaign.")
add_title(s, "Research questions")
add_body(s, [
    ("h", "RQ1 (core).  Does building the search ensemble by hazard-space coverage, rather than i.i.d. sampling from the same generator, change re-evaluated robustness?"),
    ("b", "held-out, deeply uncertain re-evaluation of the resulting Pareto-approximate policies; the single controlled contrast is fixed_probabilistic vs hazard_filling"),
    ("gap", 10),
    ("h", "RQ2.  Can FFMP re-optimization improve NYC / basin outcomes?"),
    ("b", "supply reliability · Decree flow targets · flooding · storage resilience; judged against the current FFMP policy re-evaluated on the same test ensemble"),
    ("gap", 10),
    ("h", "RQ3.  Does a variable-resolution FFMP rule structure help?"),
    ("b", "ffmp_N redraws the storage-zone discretization with N zone levels (8, 10, 12) vs the standard 24-parameter FFMP; nesting within the design campaign open (SU budget)"),
], top=1.8, size=20)

# =====================================================================
# 14 — The experiment (staged-band diagram)
# =====================================================================
s = add_slide(notes=(
    "Diagram style follows staged method figures (Bonham et al. 2024 Fig. 1; Cohen et al. 2021 "
    "Fig. 3). ONE stationary population (Kirsch-Nowak fit to the record, no climate perturbation). "
    "KEY POINT of stage 1: there is NO shared pool. Every design generates its own realizations from "
    "its own namespaced seed stream, so no two designs ever share a realization. Only hazard_filling "
    "needs a pool, and it owns it. Deep uncertainty enters ONLY in stage 4 (E_test), which spans a "
    "CMIP6-forced envelope no search ensemble contains, making re-evaluation a generalization test. "
    "Stage 3: identical N, L and NFE for both matched designs. Stage 4: one common held-out "
    "evaluation; search-time values are never compared across designs."))
add_title(s, "The planned experiment")
LX = 0.42   # stage label x
BX = 1.75   # content x start
# stage 1: GENERATE
label(s, LX, 1.72, 1.4, 0.8, "**1**\nGENERATE", size=18, color=GRAY, align=PP_ALIGN.LEFT, wrap=False)
box(s, BX, 1.5, 4.4, 0.9, "one stationary population\nKirsch–Nowak fit to the record (no climate perturbation)", size=18)
box(s, BX + 4.7, 1.5, 5.0, 0.9,
    "**each design generates its own realizations**\nfrom its own namespaced seed stream", size=18,
    line=ACCENT, line_w=1.75)
arrow(s, BX + 4.4, 1.95, BX + 4.7, 1.95)
label(s, BX + 9.85, 1.5, 1.7, 0.9, "deep uncertainty\nenters only in E_test\n(stage 4)", size=18, color=RUST, align=PP_ALIGN.LEFT)
line(s, 0.42, 2.72, 12.9, 2.72, dash="dash")
# stage 2: SELECT
label(s, LX, 3.05, 1.4, 0.8, "**2**\nBUILD", size=18, color=GRAY, align=PP_ALIGN.LEFT, wrap=False)
arrow(s, BX + 4.75, 2.4, BX + 4.75, 2.96)
chip_names = ["historic*", "fixed prob.", "hazard-fill"]
cw = 2.4
for i, nm in enumerate(chip_names):
    acc = nm.startswith("hazard-fill")
    box(s, BX + i * (cw + 0.2), 2.99, cw, 0.62, nm, size=18,
        fill=(LIGHTB if acc else WHITE), line=(ACCENT if acc else MIDGRAY), line_w=(1.75 if acc else 1.0))
label(s, BX, 3.67, 11.2, 0.4,
      "no design subsamples a shared pool; only hazard-filling needs a pool, and it owns its candidate pool + hazard image",
      size=18, color=GRAY, align=PP_ALIGN.LEFT, wrap=False)
line(s, 0.42, 4.08, 12.9, 4.08, dash="dash")
# stage 3: SEARCH
label(s, LX, 4.45, 1.4, 0.8, "**3**\nSEARCH", size=18, color=GRAY, align=PP_ALIGN.LEFT, wrap=False)
arrow(s, BX + 4.75, 3.61, BX + 4.75, 4.32)
box(s, BX, 4.35, 7.3, 0.85, "independent MM-Borg searches per design\nN = 100 × L = 10 yr, equal NFE; K draws × S seeds", size=18)
arrow(s, BX + 7.3, 4.77, BX + 7.65, 4.77)
box(s, BX + 7.65, 4.35, 3.45, 0.85, "Pareto-approximate\nsets (per design)", size=18)
line(s, 0.42, 5.35, 12.9, 5.35, dash="dash")
# stage 4: EVALUATE
label(s, LX, 5.75, 1.4, 0.8, "**4**\nEVALUATE", size=18, color=GRAY, align=PP_ALIGN.LEFT, wrap=False)
arrow(s, BX + 9.35, 5.2, BX + 4.9, 5.68)
box(s, BX, 5.65, 4.95, 0.95, "common held-out E_test (DU-forced)\nLHS over CMIP6 forcing box × realizations; never searched", size=18, line=ACCENT, line_w=1.75)
arrow(s, BX + 4.95, 6.12, BX + 5.3, 6.12)
box(s, BX + 5.3, 5.65, 5.8, 0.95, "multivariate satisficing (primary) + secondary metrics;\nnondominated sets recomputed for all designs", size=18)
label(s, BX, 6.68, 11.1, 0.62,
      "*historic = observed record, prevailing-practice reference (N = 1; cannot be size-matched); the current FFMP policy is also re-evaluated on the test ensemble as the RQ2 baseline",
      size=18, color=GRAY, align=PP_ALIGN.LEFT)

# =====================================================================
# 15 — Scenario designs table
# =====================================================================
s = add_slide(notes=(
    "The three campaign designs, from src/scenario_designs.py. ONE stationary population "
    "(Kirsch-Nowak fit to the record, no climate perturbation). Each design is built by its own "
    "published recipe from its own seed stream. fixed_probabilistic is the EXACT statistical control "
    "for hazard_filling (a uniform random size-N subset of an i.i.d. pool = N fresh i.i.d. draws), so "
    "the pair differs only in the selection rule. historic is the unmatched prevailing-practice "
    "reference (N = 1). Deep uncertainty lives only in E_test. resampled_probabilistic, "
    "input_stratified, and a DU-forced hazard_filling variant are out of the campaign (future work "
    "at most)."))
add_title(s, "Scenario designs compared (RQ1)")
rows = [
    ["Population", "Design", "Construction of the search ensemble", "Precedent"],
    ["stationary", "historic", "the observed record as one continuous trace (N = 1)", "Giuliani et al. 2016; Herman et al. 2020"],
    ["stationary", "fixed_probabilistic", "N × L realizations generated i.i.d.; frozen across the search", "Quinn et al. 2017; Zatarain Salazar et al. 2017"],
    ["stationary", "**hazard_filling**", "LHS + nearest-neighbour snap in ABSOLUTE hazard space, over its own i.i.d. stationary pool", "**proposed**"],
]
add_table(s, rows, [1.35, 2.85, 4.75, 3.2], left=0.6, top=1.6, row_h=0.8, size=18,
          highlight_row=[3], highlight_color=LIGHTB)
add_body(s, [
    ("b", "each design generates its own realizations from its own seed stream; **no design is subsampled from a shared pool**"),
    ("b", "the two matched designs run at **N = 100, L = 10 yr, equal NFE**; **fixed_probabilistic is the exact statistical control** for hazard_filling"),
    ("b", "historic is an unmatched prevailing-practice reference; deep uncertainty enters only in the held-out E_test"),
], top=4.55, size=18, space_after=4)

# =====================================================================
# 16 — Why both populations
# =====================================================================
s = add_slide(notes=(
    "The single controlled contrast rests on an exact statistical control. A uniform random size-N "
    "subset of an i.i.d. pool is distributionally identical to N fresh i.i.d. draws, so "
    "fixed_probabilistic is the EXACT control for hazard_filling: same generator, population, N, L, "
    "NFE, only the selection rule differs. This requires the pool to be sampled i.i.d. (never by "
    "LHS), enforced by an invariant test. Deep uncertainty enters only in E_test, which spans a "
    "CMIP6-forced envelope no search ensemble contains, making re-evaluation a generalization test "
    "to conditions never seen in search. This upgrades the Eker & Kwakkel (2018) null benchmark "
    "(diversity selection did not beat random) to hazard space, a real system, and draw-based "
    "replication."))
add_title(s, "The exact statistical control")
claim(s, "fixed_probabilistic is the exact statistical control for hazard_filling.")
add_body(s, [
    ("b", "same generator, same stationary population, same N, same L, same NFE — **only the selection rule differs**"),
    ("b", "so any difference in re-evaluated robustness is attributable to the **selection rule alone**; this is the Eker & Kwakkel (2018) null benchmark upgraded to hazard space, a real system, and draw-based replication"),
    ("b", "**deep uncertainty enters only in E_test**, which spans a CMIP6-forced envelope no search ensemble contains — so re-evaluation is a **generalization test** to conditions never seen in search"),
], top=2.05, size=20, space_after=9)
cb = box(s, 0.6, 5.15, 12.15, 1.55, None, fill=LIGHTB, line=ACCENT, line_w=1.25)
tf = cb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_apply_runs(tf.paragraphs[0],
            "**The exact control.**  A uniform random size-N subset of an i.i.d. pool is distributionally identical to N fresh i.i.d. draws. That is what makes fixed_probabilistic the exact statistical control for hazard_filling: they differ only in the selection rule applied to the same population law. It requires the pool to be sampled i.i.d., never by LHS — enforced by an invariant test, because nothing else would fail if it were broken.",
            20, INK)

# =====================================================================
# 17 — The controlled contrasts
# =====================================================================
s = add_slide(notes=(
    "The experimental core, in one picture. ONE controlled contrast: fixed_probabilistic -> "
    "hazard_filling. Same stationary generator, same population, same N = 100, same L = 10 yr, same "
    "NFE; only the selection RULE differs (i.i.d. sampling vs absolute hazard-space filling). Because "
    "the pool is i.i.d., fixed_probabilistic is the exact statistical control. historic is the "
    "unmatched prevailing-practice reference. RQ1: does hazard coverage beat random sampling?"))
add_title(s, "The controlled contrast")
box(s, 4.4, 2.1, 4.55, 0.95, "fixed_probabilistic\nN i.i.d. draws from the generator, frozen", size=18)
arrow(s, 6.68, 3.05, 6.68, 3.95, color=ACCENT, weight=1.9)
label(s, 6.9, 3.18, 4.0, 0.7, "only the selection\n**RULE** differs", size=18, color=ACCENT, align=PP_ALIGN.LEFT)
box(s, 4.4, 3.98, 4.55, 0.95, "hazard_filling\nLHS + snap in absolute hazard space", size=18, line=ACCENT, line_w=1.9, bold=True)
label(s, 4.4, 5.05, 4.55, 0.5, "Does hazard coverage beat\nrandom sampling?", size=20, color=GRAY)
add_body(s, [
    ("b", "the contrast holds **N = 100, L = 10 yr and NFE** fixed; only the rule that selects the N realizations changes"),
    ("b", "because the candidate pool is **i.i.d.**, fixed_probabilistic is the **exact statistical control** — a uniform random size-N subset equals N fresh draws"),
    ("b", "historic (prevailing-practice anchor, N = 1) sits outside the matched contrast; deep uncertainty enters only at re-evaluation (E_test)"),
], top=5.75, size=20, space_after=6)

# =====================================================================
# 18 — Candidate pools & forcing space
# =====================================================================
s = add_slide(notes=(
    "How the stationary candidate pool is built. The pool is sampled i.i.d., NEVER by LHS -- a "
    "random subset of an LHS design is not i.i.d., which would silently void the exact control. The "
    "pool is re-drawn on every ensemble draw: generating the pool IS part of hazard_filling's "
    "construction, so pinning it would make hazard-filling look more stable by construction. "
    "Storage: only hazard image + seeds persisted; realizations regenerate deterministically. The "
    "CMIP6 harmonic forcing box is NOT a search design; it is the envelope for the held-out E_test "
    "(shown later). Caveat: historical interannual persistence retained; claims scoped accordingly."))
add_title(s, "The candidate pool (and where the forcing box lives)")
r1y = 1.75; bh = 1.0
box(s, 0.65, r1y, 3.1, bh, "stationary Kirsch–Nowak\nfit to the record (no perturbation)", size=18)
arrow(s, 3.75, r1y + bh / 2, 4.25, r1y + bh / 2)
box(s, 4.25, r1y, 3.1, bh, "i.i.d. sampling\n(never LHS)", size=18)
arrow(s, 7.35, r1y + bh / 2, 7.85, r1y + bh / 2)
box(s, 7.85, r1y, 4.6, bh, "one i.i.d. candidate pool\n(hazard_filling selects N = 100 from it)", size=18, line=ACCENT, line_w=1.5)
add_body(s, [
    ("b", "hazard_filling **owns** its pool; the pool is **re-drawn on every ensemble draw**, so a draw re-rolls everything that is random about building the ensemble"),
    ("b", "the pool is sampled **i.i.d., never by LHS** — a random subset of an LHS design is not i.i.d., which would silently void the exact control (enforced by an invariant test)"),
    ("b", "only the **hazard image + seeds** are stored; any realization **regenerates deterministically** on demand (hundreds of GB avoided)"),
    ("b", "seed domains are **disjoint** across designs, draws, and the test ensemble, so no two ever share a realization"),
    ("b", "the **CMIP6 harmonic forcing box** (m, r₁, r₂; amplitudes sampled, phases fixed, Quinn et al. 2018) is **not a search design** — it is the envelope for the held-out E_test (Kirsch et al. 2013; Nowak et al. 2010)"),
    ("b", "historical interannual persistence is retained by the generator; claims scoped accordingly"),
], top=3.2, size=19, space_after=7)

# =====================================================================
# 19 — Hazard axes & selector
# =====================================================================
s = add_slide(notes=(
    "Selector: range-scale each hazard axis to [0,1] by its pool min-max (ABSOLUTE space, NOT "
    "empirical-CDF / rank space), draw N Latin-hypercube anchors, snap each to the nearest "
    "not-yet-selected pool member. Deterministic given the anchor seed. NO simulated annealing, no "
    "tuning, no discrepancy objective. Absolute-space filling deliberately OVER-REPRESENTS the "
    "severe (rare) hazard corners relative to their pool frequency -- that is the genuine "
    "distribution shift under study. Because the selector does not minimize a discrepancy objective, "
    "L2-star discrepancy / MST / snap distance are BUILD-QC that the intervention was administered at "
    "strength (vs a random design at the same N, m), NOT a comparison result. Axes chosen by a "
    "redundancy screen on the production pool's hazard image. A rank-space ECDF variant is a "
    "non-campaign sensitivity only."))
add_title(s, "Hazard axes & the hazard-filling selector")
add_body(s, [
    ("b", "candidate axes: SSI-6 drought-event metrics (deficit volume, duration, peak depth, onset and recovery rate) and peaks-over-threshold flood metrics (peak magnitude, pulse duration, rise rate), computed per sequence"),
    ("b", "redundancy screen retains **3–4** low-collinearity axes, balanced across the dry and wet concepts (|ρ| ≥ 0.7 clustering, Olden & Poff 2003)"),
    ("b", "selector: **range-scale each axis to [0,1] by its pool min-max (ABSOLUTE space)**, draw **N Latin-hypercube anchors**, snap each to the **nearest unused pool member**. Deterministic given the anchor seed; **no annealing, no tuning**"),
    ("b", "absolute-space filling **deliberately over-represents the severe corners** relative to their pool frequency — the distribution shift RQ1 tests (a rank-space ECDF variant is a non-campaign sensitivity only)"),
    ("b", "coverage stats (L2-star discrepancy, MST, snap distance) are **build-QC / method verification** that the intervention was administered at strength, reported against a random design at the same (N, m) — **not a comparison result**"),
], left=0.6, top=1.6, width=7.4, size=19, space_after=8)
gx, gy, gw, gh = 8.6, 2.0, 3.75, 3.75
box(s, gx, gy, gw, gh, None, fill=WHITE, line=MIDGRAY, line_w=1.0)
n_strata = 5
for k in range(1, n_strata):
    line(s, gx + gw * k / n_strata, gy, gx + gw * k / n_strata, gy + gh, dash="dash")
    line(s, gx, gy + gh * k / n_strata, gx + gw, gy + gh * k / n_strata, dash="dash")
_r.seed(5)
for _ in range(55):
    fx, fy = _r.random(), _r.random()
    dot(s, gx + 0.08 + fx * (gw - 0.2), gy + 0.08 + fy * (gh - 0.2), d=0.06, color=MIDGRAY)
# LHS anchors (hollow) snapped to selected pool members (filled)
anchors = [(0.10, 0.62), (0.31, 0.13), (0.52, 0.86), (0.70, 0.42), (0.90, 0.70),
           (0.12, 0.30), (0.51, 0.51), (0.88, 0.08), (0.30, 0.90), (0.72, 0.18)]
for i, (fx, fy) in enumerate(anchors):
    ax = gx + 0.06 + fx * (gw - 0.22)
    ay = gy + 0.06 + fy * (gh - 0.22)
    sx = ax + (0.16 if i % 2 == 0 else -0.14)
    sy = ay + (0.13 if i % 3 == 0 else -0.11)
    ring(s, ax, ay, d=0.14, color=MIDGRAY, weight=1.0)
    line(s, ax + 0.07, ay + 0.07, sx + 0.06, sy + 0.06, color=MIDGRAY, weight=0.75)
    dot(s, sx, sy, d=0.12, color=RUST)
label(s, gx, gy + gh + 0.06, gw, 0.35, "hazard axis 1 (absolute, range-scaled space)", size=18, color=GRAY)
ylab = label(s, gx - 1.05, gy + gh / 2 - 0.175, 1.6, 0.35, "hazard axis 2", size=18, color=GRAY)
ylab.rotation = 270
label(s, gx - 0.15, gy - 0.44, gw + 0.3, 0.36, "hollow = LHS anchor; filled = snapped pool member", size=18, color=RUST)

# =====================================================================
# 20 — Objectives table
# =====================================================================
s = add_slide(notes=(
    "Static 1954 Decree goalposts, never the live FFMP step-down targets (a policy must not lower "
    "its own goalpost). Stable tail/count forms replace worst-case extremes. Trenton stands in for "
    "the salt front (physically redundant target, cleaner signal). Full formulations and rationale "
    "per objective in supplement."))
add_title(s, "Objectives: seven, spanning basin stakeholders")
rows = [
    ["", "Objective", "Temporal metric (single trace)", "Dir"],
    ["1", "NYC delivery reliability", "fraction of weeks delivery ≥ 99% of Decree-capped demand", "max"],
    ["2", "NYC delivery deficit", "CVaR₉₀ of weekly deficit (% of 800 MGD)", "min"],
    ["3", "Montague flow reliability", "fraction of weeks ≥ 1,750 cfs Decree target", "max"],
    ["4", "Montague flow deficit", "CVaR₉₀ of weekly deficit (% of target)", "min"],
    ["5", "Trenton flow reliability", "fraction of weeks ≥ Decree target", "max"],
    ["6", "Downstream flood days", "days any tail gauge ≥ NWS minor flood stage", "min"],
    ["7", "NYC storage resilience", "5th percentile of daily storage (% capacity)", "max"],
]
add_table(s, rows, [0.5, 3.5, 7.15, 0.95], left=0.6, top=1.55, row_h=0.46, size=18)
add_body(s, [
    ("b", "goalposts are the **static 1954 Decree quantities**, never the live FFMP step-down targets"),
    ("b", "reliability, tail and count forms follow Hashimoto et al. (1982), Rockafellar & Uryasev (2000), Quinn et al. (2017); formulations in supplement"),
], top=5.45, size=19, space_after=6)

# =====================================================================
# 21 — Two-layer aggregation
# =====================================================================
s = add_slide(notes=(
    "How a (realizations x time) simulation becomes one scalar per objective during ensemble "
    "search. Stage 1: annual metric on every (realization x water-year) unit. Stage 2: "
    "per-objective operator over the pooled N*L unit-years. Every operator follows published "
    "search-time practice (WaterPaths-lineage frequency; Quinn's percentile over annual units). "
    "Identical scheme and unit denominator for every ensemble design; historic keeps single-trace "
    "metrics. Failure criteria and epsilons from the sensitivity experiment (in progress)."))
add_title(s, "Scoring a policy on an ensemble: two-layer aggregation")
gx, gy = 0.7, 1.9
cw, ch = 0.34, 0.30
n_rows_g, n_cols_g = 5, 8
for i in range(n_rows_g):
    for j in range(n_cols_g):
        c = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(gx + j * cw), Inches(gy + i * ch),
                               Inches(cw - 0.04), Inches(ch - 0.04))
        c.fill.solid(); c.fill.fore_color.rgb = LIGHTB
        c.line.color.rgb = MIDGRAY; c.line.width = Pt(0.5); c.shadow.inherit = False
hl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(gx + 3 * cw), Inches(gy + 2 * ch),
                        Inches(cw - 0.04), Inches(ch - 0.04))
hl.fill.solid(); hl.fill.fore_color.rgb = RUST; hl.line.fill.background(); hl.shadow.inherit = False
label(s, gx - 0.12, gy + n_rows_g * ch + 0.05, n_cols_g * cw + 0.3, 0.35, "L water-years", size=18, color=GRAY, align=PP_ALIGN.CENTER)
label(s, gx - 0.55, gy + 0.35, 0.5, 1.2, "N", size=18, color=GRAY)
ax0 = gx + n_cols_g * cw + 0.15
arrow(s, ax0, gy + 0.75, ax0 + 0.45, gy + 0.75)
box(s, ax0 + 0.5, gy + 0.28, 3.15, 0.95, "**stage 1**: annual metric per\n(realization × year) unit", size=18)
arrow(s, ax0 + 3.65, gy + 0.75, ax0 + 4.1, gy + 0.75)
box(s, ax0 + 4.15, gy + 0.28, 3.15, 0.95, "**stage 2**: unit operator over\npooled N·L unit-years", size=18)
arrow(s, ax0 + 7.3, gy + 0.75, ax0 + 7.7, gy + 0.75)
box(s, ax0 + 7.75, gy + 0.28, 0.75, 0.95, "fᵢ", size=20, font=TITLE_FONT)
add_body(s, [
    ("b", "reliability objectives: **frequency of non-failure years** (Zeff et al. 2014; Trindade et al. 2017; Gold et al. 2023)"),
    ("b", "deficit / storage objectives: **worst 1st-percentile unit-year** (Quinn et al. 2017, 2018)"),
    ("b", "flood days: **mean annual count**, P99 variant registered (Quinn et al. 2017)"),
    ("b", "identical scheme and unit denominator across designs (two-layer vocabulary, Hamilton et al. 2022); failure criteria and ε from the sensitivity experiment"),
], top=3.85, size=20)

# =====================================================================
# 22 — Controls
# =====================================================================
s = add_slide(notes=(
    "The budget control is a CONSEQUENCE of the sizing choice, not something imposed. Because N and "
    "L are common to every matched design, per-evaluation cost, warm-up, scenario-years and "
    "wall-clock are identical, so equal-NFE and equal-scenario-years coincide: ONE budget "
    "condition, no arms, no composition-vs-search-effort confound. A common (N, L) is REQUIRED, not "
    "convenient: with different L the selection rule would be confounded with record length. N is "
    "bounded below by the fill requirement. Replication: a draw re-rolls EVERYTHING, including the "
    "pool -- otherwise a hazard-filling draw would vary only its anchor plan while a "
    "fixed_probabilistic draw re-rolls its whole sample, and hazard-filling would look more stable "
    "by construction."))
add_title(s, "Controls for a fair comparison")
claim(s, "N = 100 and L = 10 yr for every matched design → 1,000 scenario-years per evaluation, at equal NFE.")
add_body(s, [
    ("b", "because N and L are common, per-evaluation cost, warm-up, scenario-years and wall-clock are **identical** — equal-NFE and equal-scenario-years **coincide**. One budget condition, no arms, no confound between composition and search effort"),
    ("b", "the common (N, L) is **required, not convenient**: if L differed across designs, the selection rule would be confounded with record length"),
    ("b", "N is bounded below by the **fill requirement** — at m = 4 hazard axes, N = 100 gives ~3.2 points per dimension (~4.6 at m = 3), the smallest defensible fill. This is why long records are not viable: at a fixed per-evaluation budget L = 50 forces N ≈ 20, and space-filling in 4-D with 20 points is noise"),
    ("b", "**NFE = 500,000 per search is a target** (revisable once initial searches reveal convergence; the runtime archive keeps intermediate NFE so a lower budget can be justified after the fact)"),
    ("b", "**replication (targets)**: K = 3 ensemble draws × S = 2 seeds per matched design (historic K = 1); revisable from a pilot minimum-detectable-effect calc. A **draw** = the construction re-run from scratch with a fresh seed, **including its pool**, and is the unit of analysis (effective n ≈ K)"),
    ("b", "search-time objective values are **never compared across designs**; convergence is a per-design diagnostic only"),
], top=2.15, size=19, space_after=9)

# =====================================================================
# 23 — Comparison metrics
# =====================================================================
s = add_slide(notes=(
    "The finalized metric set (src/robustness.py). PRIMARY: multivariate Starr (1962) domain "
    "criterion -- realization-level, the standard unit of the Herman/Trindade/Gold lineage, and it "
    "converges at 50-300 scenarios (Bonham 2024). Everything else is FREE from the same persisted "
    "(solution x realization x objective) cube. Three deliberate exclusions, each with a reason a "
    "reviewer will accept. NO perfect-foresight optimization appears anywhere in the study."))
add_title(s, "Comparison metrics: the held-out metric set")
pb = box(s, 0.6, 1.42, 12.15, 0.95, None, fill=LIGHTB, line=ACCENT, line_w=1.25)
tf = pb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_apply_runs(tf.paragraphs[0],
            "**Primary: multivariate satisficing (Starr 1962 domain criterion)** — the fraction of test realizations in which a policy meets ALL criteria jointly. The realization is the unit.",
            20, INK)
add_body(s, [
    ("b", "the standard robustness unit of the Herman (2014, 2015) / Trindade (2017) / Gold (2022, 2023) lineage; converges at 50–300 scenarios (Bonham et al. 2024)"),
    ("b", "**run-level endpoint (pre-specified)**: the **max satisficing fraction** attained in a run's re-evaluated set, with a **leave-one-out reference-set correction** and **robustness-space hypervolume** as bounding co-reports"),
    ("b", "**secondary, all free from the same persisted cube**: univariate satisficing; **Laplace / mean** and **maximin** (risk-neutral and risk-averse anchors — a single robustness family is never sufficient, Herman et al. 2015; McPhail et al. 2018); **improvement over the status quo** (shortfall vs the current FFMP policy on the same test ensemble — the **only** fixed-reference regret-type quantity, design-independent, needs no optimization, Kasprzyk et al. 2013)"),
    ("b", "**attainability screen** (free): which test realizations no policy from any design can win — separating *this design searched badly* from *this scenario is impossible for anyone* (Shavazipour et al. 2021 found 23% unwinnable)"),
    ("b", "ranking agreement via **Kendall's τ_b, computed across the design rankings**"),
], top=2.55, size=19, space_after=7)
xb = box(s, 0.6, 5.5, 12.15, 1.55, None, fill=WHITE, line=RUST, line_w=1.25)
tf = xb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p0 = tf.paragraphs[0]
_apply_runs(p0, "**Deliberately excluded.**  **Regret** — Cohen et al. (2021) baseline regret needs a perfect-foresight MOEA run per scenario (97 optimizations, ~3,233 CPU-h) and does not scale to a 10⁵–10⁶ pool; regret-from-best is set-relative and design-coupled (dropping one design changes every other design's score) and never converges on a tail objective (Bonham et al. 2024). **No perfect-foresight optimization appears anywhere in this study.**", 18, INK)
p1 = tf.add_paragraph()
_apply_runs(p1, "**Pooled-reference-set hypervolume across designs** — contributor bias, cardinality asymmetry, and noise-induced spurious dominance (a noisier design is flattered). Reference sets are built per design; cross-design comparison is re-evaluation only (Zatarain Salazar et al. 2017, Section 5.3).", 18, INK)

# =====================================================================
# 24 — Threshold sweep
# =====================================================================
s = add_slide(notes=(
    "A MAIN-TEXT figure, not a robustness check. Quinn et al. 2020: robustness-rank agreement "
    "ACROSS scenario designs degrades as the satisficing criterion becomes more stringent -- 'the "
    "more conservative one wants to be in finding robust policies, the harder it is to choose this "
    "consistently across experimental designs.' So the design effect is largest at the conservative "
    "end and a single fixed threshold could manufacture or hide the entire result. We sweep a "
    "stringency grid and ask whether the DESIGN RANKING is threshold-invariant. Free: the raw cube "
    "is re-scored offline, no re-simulation. Nobody in this lineage has done this."))
add_title(s, "Satisficing thresholds are conventions, so we sweep them")
claim(s, "Rank agreement across scenario designs degrades as the criterion tightens (Quinn et al. 2020).")
add_body(s, [
    ("b", "nobody in this lineage derives a threshold: Zeff et al. (2014) elicited them, and every later number descends by convention and drifts without stated reason (worst-case cost 5% → 10%; restriction frequency 20% → 10%; reliability 99% → 98.5% because nothing met 99%)"),
    ("b", "so the **design effect is largest at the conservative end**, and a single fixed threshold could **manufacture or hide the entire result**"),
    ("b", "we anchor each criterion on a **Decree / FFMP goalpost** where one exists, and sweep the rest across a stringency grid"),
    ("b", "the reported quantity is whether the **design ranking** — not the robustness value — is threshold-invariant (Kendall's τ_b vs the default): a **main-text figure**, not a robustness check"),
    ("b", "free: the persisted (solution × realization × objective) cube is re-scored offline; nothing is re-simulated"),
], left=0.6, top=2.1, width=7.6, size=19, space_after=9)
gx, gy, gw, gh = 8.6, 2.4, 3.9, 3.3
box(s, gx, gy, gw, gh, None, fill=WHITE, line=MIDGRAY, line_w=1.0)
arrow(s, gx + 0.5, gy + gh - 0.5, gx + 0.5, gy + 0.25, color=MIDGRAY, weight=1.2)
arrow(s, gx + 0.5, gy + gh - 0.5, gx + gw - 0.2, gy + gh - 0.5, color=MIDGRAY, weight=1.2)
series = [
    ([(0.05, 0.30), (0.40, 0.55), (0.75, 1.05), (1.00, 1.55)], MIDGRAY, 1.5),
    ([(0.05, 0.55), (0.40, 0.85), (0.75, 1.50), (1.00, 2.05)], MIDGRAY, 1.5),
    ([(0.05, 0.80), (0.40, 1.20), (0.75, 1.35), (1.00, 1.35)], MIDGRAY, 1.5),
    ([(0.05, 0.45), (0.40, 0.62), (0.75, 0.72), (1.00, 0.80)], RUST, 2.4),
]
for pts_s, col, wgt in series:
    for k in range(len(pts_s) - 1):
        x1 = gx + 0.7 + pts_s[k][0] * 2.6;   y1 = gy + 0.3 + pts_s[k][1]
        x2 = gx + 0.7 + pts_s[k + 1][0] * 2.6; y2 = gy + 0.3 + pts_s[k + 1][1]
        line(s, x1, y1, x2, y2, color=col, weight=wgt)
label(s, gx + 0.55, gy + gh - 0.42, gw - 0.75, 0.32, "criterion stringency →", size=18, color=GRAY)
ylab = label(s, gx - 0.85, gy + 1.35, 1.9, 0.35, "robustness", size=18, color=GRAY)
ylab.rotation = 270
label(s, gx + 1.9, gy + 0.15, 1.9, 0.32, "hazard-fill", size=18, color=RUST)
label(s, gx - 0.15, gy + gh + 0.06, gw + 0.3, 0.6,
      "conceptual: rankings diverge at the tight end;\nτ_b of the design ranking reported at every grid point",
      size=18, color=GRAY)

# =====================================================================
# 25 — Mechanism test
# =====================================================================
s = add_slide(notes=(
    "Scenario discovery in HAZARD space is the MECHANISM TEST, not decoration. Boosted trees (Gold "
    "2022/2023 configuration) on each design's held-out failures, labelled by the joint satisficing "
    "criterion, in hazard space. The falsifiable prediction: if hazard-filling works, its policies "
    "should fail in the hazard region that the design UNDER-COVERED. Caveat to state up front: the "
    "raw coverage-deficit -> failure AUC is inflated by GEOMETRY -- nearest-neighbour distance grows "
    "toward the manifold boundary -- so the verdict is read off AUC minus a random-ensemble null, "
    "never raw AUC."))
add_title(s, "Scenario discovery in hazard space: the mechanism test")
claim(s, "If hazard-filling works, its policies should fail where the design under-covered.")
add_body(s, [
    ("b", "boosted trees (Gold et al. 2022, 2023 configuration) fit to each design's **held-out failure realizations**, labelled by the **joint** satisficing criterion, **in hazard space**"),
    ("b", "this is the **falsifiable prediction** of the coverage claim, and strictly stronger than a correlational coverage → robustness association"),
    ("b", "**caveat, stated up front**: the raw coverage-deficit → failure AUC is inflated by **geometry** (nearest-neighbour distance grows toward the manifold boundary), so the verdict is read off **AUC minus a random-ensemble null**, never raw AUC"),
    ("b", "correlated hazard axes destabilize factor importances (Quinn et al. 2020), so the axes are redundancy-screened before this is run"),
], left=0.6, top=2.1, width=7.5, size=19, space_after=10)
gx, gy, gw, gh = 8.5, 2.35, 4.05, 3.5
box(s, gx, gy, gw, gh, None, fill=WHITE, line=MIDGRAY, line_w=1.0)
und = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(gx + 2.25), Inches(gy + 0.25),
                         Inches(1.6), Inches(1.5))
und.fill.solid(); und.fill.fore_color.rgb = LIGHTB
und.line.color.rgb = ACCENT; und.line.width = Pt(1.0); und.shadow.inherit = False
ln_el = und.line._get_or_add_ln(); d_el = ln_el.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln_el.append(d_el)
_r.seed(31)
for _ in range(30):
    dot(s, gx + 0.2 + _r.random() * (gw - 0.5), gy + 0.2 + _r.random() * (gh - 0.7), d=0.07, color=MIDGRAY)
for _ in range(14):
    dot(s, gx + 2.32 + _r.random() * 1.35, gy + 0.32 + _r.random() * 1.25, d=0.1, color=RUST)
label(s, gx + 1.9, gy + 1.82, 2.3, 0.32, "failures", size=18, color=RUST)
label(s, gx + 1.6, gy - 0.42, 2.6, 0.36, "under-covered region", size=18, color=ACCENT)
label(s, gx, gy + gh + 0.06, gw, 0.32, "hazard axis 1", size=18, color=GRAY)
ylab = label(s, gx - 0.9, gy + 1.4, 1.9, 0.35, "hazard axis 2", size=18, color=GRAY)
ylab.rotation = 270

# =====================================================================
# 26 — Threats to validity
# =====================================================================
s = add_slide(notes=(
    "Name them before a reviewer does. Eker & Kwakkel 2018 is a NULL RESULT and our "
    "fixed_probabilistic -> hazard_filling contrast IS that benchmark. Giuliani & "
    "Castelletti 2016 Fig. 4 predicts hazard-filling will be PENALIZED for search-measure mismatch "
    "-- that is the null this study must beat, and we say so. Bartholomew & Kwakkel 2020 already "
    "wrote the reviewer's objection: selection bias is NOT corrected by held-out re-evaluation, and "
    "it is precisely what the experiment measures."))
add_title(s, "Threats to validity, named before a reviewer names them")
add_body(s, [
    ("h", "Eker & Kwakkel (2018) is a null result"),
    ("b", "diversity-based scenario selection did **not** beat random selection. Our fixed_probabilistic vs hazard_filling contrast **is** that benchmark. Differentiators: their diversity is in *outcome* space on the Lake Problem (little scenario-to-outcome leverage), ours is in absolute hazard space on a real system; and our statistic has far more power (K draws × S seeds, the draw as the unit of analysis, vs counting solutions above a group median)"),
    ("h", "Giuliani & Castelletti (2016): search-measure mismatch is a one-directional penalty"),
    ("b", "policies designed under one aggregation rule and scored under another are **dominated** (their Fig. 4). Hazard-filling searches under a distorted measure and is scored under the test measure, so Giuliani predicts it will be **penalized**. **That is the null this study must beat**"),
    ("h", "Bartholomew & Kwakkel (2020) wrote the objection already"),
    ("b", "selecting scenarios from a chosen region “intrinsically biases subsequent results towards solutions that do well in this region … no a-priori reason to assume [they] might not be vulnerable in a different way.” Our answer: the distortion is deliberate, stated and coverage-motivated; held-out re-evaluation corrects **evaluation** bias; **selection** bias is not corrected, and is precisely what the experiment measures"),
    ("h", "Degeneracy"),
    ("b", "a metric can be stable, optimizable and still perverse (Huang et al. 2025: a deviation metric is minimized by being uniformly terrible), so **raw performance distributions are co-reported** with every robustness number"),
], top=1.5, size=18, space_after=4)

# =====================================================================
# 27 — Evaluation funnel
# =====================================================================
s = add_slide(notes=(
    "Single point of cross-design comparison, and the SOLE carrier of deep uncertainty. E_test is an "
    "LHS over the CMIP6 harmonic forcing box crossed with many realizations per point -- LHS, NOT "
    "i.i.d. -- spanning a forced envelope no search ensemble contains. So re-evaluation is a "
    "GENERALIZATION test to conditions never seen in search, and E_test is structurally distinct "
    "from both designs (it favours neither). No robustness number is an expectation: a satisficing "
    "fraction over E_test is a coverage-weighted count over a designed exploration of the deeply "
    "uncertain box; the comparison is valid because E_test is identical across designs, not because "
    "it is probability-faithful. An optional 2nd construction (multi-site HMM) would test ranking "
    "stability across generator families. Nondominated sets recomputed from re-evaluated values for "
    "all designs alike."))
add_title(s, "Evaluation: one common held-out test ensemble")
names = ["historic", "fixed prob.", "hazard-fill", "current FFMP"]
by0 = 2.05
for i, nm in enumerate(names):
    yy = by0 + i * 0.72
    base = (nm == "current FFMP")
    box(s, 0.7, yy, 2.9, 0.55, ("baseline: " if base else "Pareto set: ") + nm, size=18,
        fill=WHITE, line=MIDGRAY, line_w=1.0)
    arrow(s, 3.6, yy + 0.27, 4.75, 3.95, color=MIDGRAY, weight=1.0)
box(s, 4.8, 3.1, 3.6, 1.7, "**E_test**, held-out (DU-forced)\nLHS over the CMIP6 forcing box × realizations\nspans a forced envelope no search ensemble contains", size=18, line=ACCENT, line_w=1.75)
arrow(s, 8.4, 3.95, 9.0, 3.95)
box(s, 9.0, 3.05, 3.7, 1.8, "**multivariate satisficing** (primary)\n+ Laplace · maximin · vs status quo\n+ attainability screen", size=18)
add_body(s, [
    ("b", "E_test is the **sole carrier of deep uncertainty** — search is stationary, so re-evaluation is a **generalization test** to conditions never seen in search; E_test is structurally distinct from both designs and favours neither"),
    ("b", "**no robustness number is an expectation** — a satisficing fraction over E_test is a coverage-weighted count; the comparison is valid because E_test is **identical across designs**, not because it is probability-faithful (Quinn et al. 2020)"),
    ("b", "E_test is **LHS, not i.i.d.**, and **never the source of any search ensemble**; seed domains are disjoint by construction. Optional 2nd construction (multi-site HMM) tests ranking stability across generator families"),
], top=6.05, size=18, space_after=4)

# =====================================================================
# 28 — Status & open decisions
# =====================================================================
s = add_slide(notes=(
    "Left: verified as of July 2026. Right: open decisions in priority order, where committee "
    "input is most valuable. SU table lands when the Anvil scaling experiment finishes."))
add_title(s, "Status & open decisions")
add_body(s, [
    ("h", "Built & verified"),
    ("b", "end-to-end pipeline at test scale"),
    ("b", "per-design generation, disjoint seed domains"),
    ("b", "deterministic pools + on-demand regeneration"),
    ("b", "hazard-filling selector (LHS + nearest-neighbour, absolute space) wired"),
    ("b", "Anvil scaling experiment running"),
    ("gap", 8),
    ("h", "Next"),
    ("b", "production candidate pools, then the axis screen, then the sensitivity experiments, then the campaign"),
], left=0.6, top=1.7, width=5.9, size=19)
ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.8), Inches(1.7), Inches(6.8), Inches(6.6))
ln.line.color.rgb = MIDGRAY; ln.line.width = Pt(1.0); ln.shadow.inherit = False
add_body(s, [
    ("h", "Open, input requested"),
    ("b", "**coverage adequacy at m ≈ 4, N = 100** — uniform filling is sparse; must be **demonstrated** against a random-design null, not asserted (the biggest methodological risk)"),
    ("b", "NFE (target, revisable); K / S replication targets"),
    ("b", "pool size P; the hazard-axis set (post-screen)"),
    ("b", "E_test sizing; thresholds / ε; whether to stand up the optional 2nd (HMM) E_test"),
    ("b", "which design carries the RQ3 variable-resolution (ffmp_N) sweep"),
], left=7.2, top=1.7, width=5.5, size=19)

# =====================================================================
# SUPPLEMENTAL
# =====================================================================
s = add_slide(notes="Supplemental divider.")
divider(s, [("Supplemental", True)])

s = add_slide(notes="Plain-language building blocks for the objective slides.")
add_title(s, "Notation: building blocks")
add_body(s, [
    ("f", "shortfall:  (gap)₊ = max(gap, 0)"),
    ("f", "unit-year = one water year (Oct 1–Sep 30);  an L-yr record scores L−1 unit-years"),
    ("f", "CVaR₉₀(series) = mean of the worst 10% of weekly values"),
    ("f", "P₉₉ / P₀₁(units) = worst / 1st-percentile unit-year"),
], top=1.5, size=20, space_after=8)
add_body(s, [
    ("b", "flows / demands in MGD; storage in MG; first 365 days (warm-up) dropped"),
    ("b", "Decree goalposts: NYC 800 MGD; Montague 1,131.05 MGD (1,750 cfs); Trenton 1,938.95 MGD; NYC capacity 270,837 MG"),
    ("b", "every objective is two-layer — an annual metric per unit-year, then a unit operator over the pooled unit-years. The SAME function for every design (historic = 1 record × 76 unit-years; ensembles = N × (L−1)), so only the scenario set varies"),
], top=3.9, size=19)


def objective_slide(num, name, direction, formulas, bullets, notes):
    s = add_slide(notes=notes)
    add_title(s, "Objective %d: %s  (%s)" % (num, name, direction))
    add_body(s, [("f", f) for f in formulas], top=1.5, size=20, space_after=6)
    items = [("b", b) for b in bullets]
    add_body(s, items, top=1.5 + 0.5 * len(formulas) + 0.3, size=19)
    return s

objective_slide(
    1, "NYC delivery reliability", "higher is better",
    ["entitlement:  Eₜ = min(demandₜ, allowanceₜ)      [allowance = running-avg bank]",
     "failing week:  Σ delivery < 99% × Σ E       failure-year:  ≥ k failing weeks  (k = 1)",
     "f₁ = (unit-years without failure) ÷ (total unit-years)"],
    ["the Decree caps the running AVERAGE diversion, not any single day — demand is never clipped at a flat 800 MGD, and a spike is owed whenever prior under-use has banked the allowance for it",
     "the bank accrues 800 − delivery per day (reset Jun 1) at the STATIC Decree right, so a policy's own drought step-downs cannot lower its own goalpost",
     "failure-year frequency is the citable search-time satisficing operator (Zeff et al. 2014; Trindade et al. 2017)",
     "re-evaluation metric: whole-record weekly reliability, per realization"],
    "Annual failure-year frequency vs the running-average entitlement; same objective for every design.")

objective_slide(
    2, "NYC delivery deficit", "lower is better",
    ["weekly deficit = (E − delivery)₊ ÷ 800 MGD × 100%",
     "annual metric = CVaR₉₀ of that unit-year's weekly deficits",
     "f₂ = P₉₉ across pooled unit-years   (worst 1st-percentile unit-year)"],
    ["magnitude of shortfalls below the running-average entitlement E, not just their frequency; scaled to the fixed 800 MGD right",
     "the tail average is far stabler than the single worst week (Rockafellar & Uryasev 2000; Quinn et al. 2017)",
     "re-evaluation metric: whole-record CVaR₉₀, per realization"],
    "Within-year CVaR90, then worst-1%-unit-year across the pool.")

objective_slide(
    3, "Montague flow reliability", "higher is better",
    ["failing week:  mean flow < 1,131.05 MGD      [1,750 cfs]",
     "f₃ = (unit-years with < k failing weeks) ÷ (total unit-years)"],
    ["vs the fixed Decree target, never the live FFMP step-down target",
     "won't reach 100%: FFMP drought rules cut releases below target by design, keeping the signal continuous",
     "failure-year frequency (Zeff et al. 2014; Trindade et al. 2017)",
     "re-evaluation metric: whole-record weekly reliability, per realization"],
    "Annual failure-year frequency vs static Decree target.")

objective_slide(
    4, "Montague flow deficit", "lower is better",
    ["weekly deficit = (1,131.05 − flow)₊ ÷ 1,131.05 × 100%",
     "annual metric = CVaR₉₀ of that unit-year's weekly deficits",
     "f₄ = P₉₉ across pooled unit-years"],
    ["depth of target misses in the bad weeks, as % of target",
     "Montague flow is storm-dominated, so the single worst week is mostly exogenous noise; use the tail average (Quinn et al. 2017)",
     "re-evaluation metric: whole-record CVaR₉₀, per realization"],
    "Within-year CVaR90, then worst-1%-unit-year.")

objective_slide(
    5, "Trenton flow reliability", "higher is better",
    ["failing week:  mean flow < 1,938.95 MGD",
     "f₅ = (unit-years with < k failing weeks) ÷ (total unit-years)"],
    ["lower-basin / NJ obligation; direct representation of a co-equal Decree party (Trindade et al. 2017; Hadjimichael et al. 2020)",
     "replaces the salt-front objective: physically redundant target, cleaner signal in drought",
     "re-evaluation metric: whole-record weekly reliability, per realization"],
    "Trenton reliability replacing salt front.")

objective_slide(
    6, "Downstream flood days", "lower is better",
    ["annual metric = days any tail gauge ≥ NWS minor flood stage, per unit-year",
     "f₆ = mean across pooled unit-years   (expected annual flood days)",
     "gauges: Hale Eddy, Fishs Eddy, Bridgeville"],
    ["tail gauges respond to release decisions, so flooding here is attributable to operations",
     "count-over-threshold avoids the expected-damage trap (Quinn et al. 2017); a P₉₉ variant is registered because an expectation can mask floods",
     "re-evaluation metric: whole-record flood-day count, per realization"],
    "Expected annual flood days at minor stage.")

objective_slide(
    7, "NYC storage resilience", "higher is better",
    ["storage% = combined NYC storage ÷ 270,837 MG × 100%",
     "annual metric = that unit-year's minimum daily storage%",
     "f₇ = P₀₁ across pooled unit-years   (1st-percentile unit-year)"],
    ["how depleted the system routinely runs; a low percentile is stabler than the single-day minimum, which one event dominates (Quinn et al. 2017)",
     "re-evaluation metric: whole-record P₅ of daily storage%, per realization"],
    "Annual minimum storage, then P01 across the pool.")

# ---- Forcing supplement ----
s = add_slide(notes=(
    "Each CMIP6 (GCM x SSP x period) run gives a 12-month change-factor profile; the log profile "
    "is decomposed into a low-order harmonic with phases FIXED at the canonical CMIP6 seasonal "
    "shape; only amplitudes sampled. Left: best fits; right: worst (flat profiles, low R2 but "
    "small amplitude, harmless)."))
add_title(s, "Supplemental: fixed-phase harmonic fits to CMIP6 (E_test forcing)")
pic(s, "image16.png", 1.15, 1.65, w=11.0)
add_body(s, [
    ("b", "sampled axes: volume m · seasonal amplitude r₁ · shoulder shape r₂; phases fixed at the canonical CMIP6 shape (Quinn et al. 2018). This parameterization defines the **held-out E_test forcing envelope**, not a search design"),
], top=6.3, size=19, space_after=6)

s = add_slide(notes=(
    "The DMDU box = empirical 90% range of the CMIP6 fitted parameters (markers). This box is the "
    "forcing envelope for the held-out E_test: E_test is an LHS over this box crossed with "
    "realizations. It is NOT a search design (search is stationary; no input_stratified design). "
    "Panel c: phase held at the canonical winter peak."))
add_title(s, "Supplemental: the CMIP6 90% amplitude envelope (E_test)")
pic(s, "image18.png", 0.75, 1.6, w=11.8)
add_body(s, [
    ("b", "the **E_test forcing envelope**: deeply uncertain magnitude combinations the GCMs did not jointly produce, with a physically plausible seasonal shape; sampled by LHS to build E_test"),
], top=6.5, size=19, space_after=6)

s = add_slide(notes=(
    "Sanity check: the sampled change factors and the resulting monthly NYC inflows span the raw "
    "CMIP6 monthly range."))
add_title(s, "Supplemental: sampled forcing spans the CMIP6 monthly range")
pic(s, "image19.png", 0.9, 1.7, w=11.6)

s = add_slide(notes=(
    "Scaling experiment run on Anvil; sets production parallel geometry and the SU budget that gates "
    "campaign scope. Measured cost surface: 173.8 s/eval at N=100, L=10 on the trimmed model; full "
    "model 1.16x; ~33,300 SU per 500k-NFE search; full campaign ~415,000 of the 750,000-SU "
    "allocation; re-evaluation nearly free."))
add_title(s, "Supplemental: Anvil scaling & measured cost surface")
add_body(s, [
    ("h", "Parallelism levers"),
    ("b", "nodes · cores per node · MM-Borg islands · workers per island · realizations per simulation"),
    ("gap", 6),
    ("h", "Measured cost (Anvil)"),
    ("b", "**173.8 s per evaluation** (N = 100, L = 10, trimmed model); full model **1.16×**"),
    ("b", "**~33,300 SU per 500k-NFE search**; full campaign **~415,000 of the 750,000-SU allocation**; re-evaluation nearly free"),
    ("b", "packing sweep + Borg strong scaling set the production parallel geometry"),
], top=1.7, size=20)

# =====================================================================
# PLANNED MANUSCRIPT FIGURE SEQUENCE
# =====================================================================
s = add_slide(notes=(
    "Conceptual figure sequence for the manuscript: guides narrative and the processing/metrics "
    "to build. Inspired by the figure sequences of Cohen et al. 2021 (Figs. 2-6), Trindade et al. "
    "2017 (Figs. 5-8), Bonham et al. 2024 (Fig. 1). Details finalized after the sensitivity "
    "experiments."))
divider(s, [("Planned Manuscript Figure Sequence", True), ("(conceptual; details set after the sensitivity experiments)", False)])


def mock_frame(s, x=7.05, y=2.0, w=5.65, h=4.5):
    label(s, x, y - 0.42, w, 0.34, "conceptual mockup", size=18, color=MIDGRAY, italic=True)
    box(s, x, y, w, h, None, fill=WHITE, line=MIDGRAY, line_w=1.0)
    return x, y, w, h


def fig_slide(num, short, bullets, notes):
    s = add_slide(notes=notes)
    add_title(s, "Planned Fig. %d: %s" % (num, short))
    add_body(s, bullets, left=0.6, top=1.6, width=6.2, size=20, space_after=9)
    return s


def polyline(s, pts, color, weight=1.5):
    for k in range(len(pts) - 1):
        line(s, pts[k][0], pts[k][1], pts[k + 1][0], pts[k + 1][1], color=color, weight=weight)


# ---- Overview list ----
s = add_slide(notes=(
    "One-line map of the planned figures. Figs 1-2 methods; 3-7 RQ1 (headline, threshold "
    "stringency, the controlled contrast, policy differences, mechanism); 8 RQ2; 9 RQ3. There is "
    "no robustness-vs-regret figure (regret is not computed) and no budget-interaction figure "
    "(there is one budget condition). Leave-one-out reference-set hypervolume stays in the "
    "supplement, per Zatarain Salazar 2017. Three designs, one stationary population."))
add_title(s, "Planned figure sequence (conceptual)")
add_body(s, [
    ("b", "**Fig. 1**  System & experiment: the NYC-DRB testbed and the four-stage comparison workflow."),
    ("b", "**Fig. 2**  Search ensembles in hazard space: what each of the three designs samples from the stationary population."),
    ("b", "**Fig. 3**  Robustness by design: held-out multivariate satisficing per design and objective (RQ1 headline)."),
    ("b", "**Fig. 4**  Threshold sweep: design robustness vs criterion stringency, with τ_b of the design ranking (the credibility figure)."),
    ("b", "**Fig. 5**  The controlled contrast: fixed_probabilistic vs hazard_filling robustness, with historic as reference."),
    ("b", "**Fig. 6**  Policy differences: do designs change the policies, not just the scores (RQ1)."),
    ("b", "**Fig. 7**  Mechanism test: failure rate vs coverage deficit, against the random-ensemble null (RQ1)."),
    ("b", "**Fig. 8**  Re-optimized vs current FFMP: objective gains and drought behavior (RQ2)."),
    ("b", "**Fig. 9**  Rule resolution: robustness vs FFMP zone count, ffmp_N (RQ3)."),
], top=1.55, size=20, space_after=7)

# ---- Fig 1: system & experiment ----
s = fig_slide(1, "system & experimental design",
    [("b", "shows: the NYC-DRB system (reservoirs, diversion, Decree targets) and the four-stage comparison workflow in one methods figure"),
     ("b", "role: orients the reader; every results figure refers back to its stages"),
     ("gap", 2),
     ("c", "style cue: staged workflow figures of Bonham et al. 2024 (Fig. 1) and Cohen et al. 2021 (Fig. 3).")],
    ("Panel (a): basin schematic (could be a proper map). Panel (b): the 4-stage GENERATE/BUILD/"
     "SEARCH/EVALUATE workflow. Processing needs: none beyond existing diagrams."))
fx, fy, fw, fh = mock_frame(s)
label(s, fx + 0.1, fy + 0.1, 2.4, 0.3, "(a) system", size=18, color=GRAY, align=PP_ALIGN.LEFT)
box(s, fx + 0.35, fy + 0.55, 0.62, 0.34, None, fill=LIGHTB, line=MIDGRAY)
box(s, fx + 1.1, fy + 0.55, 0.62, 0.34, None, fill=LIGHTB, line=MIDGRAY)
box(s, fx + 1.85, fy + 0.55, 0.62, 0.34, None, fill=LIGHTB, line=MIDGRAY)
arrow(s, fx + 1.4, fy + 0.95, fx + 1.4, fy + 1.5, color=MIDGRAY, weight=1.0)
box(s, fx + 1.0, fy + 1.55, 0.85, 0.34, None, fill=LIGHTB, line=MIDGRAY)
arrow(s, fx + 1.4, fy + 1.9, fx + 1.4, fy + 2.45, color=MIDGRAY, weight=1.0)
box(s, fx + 1.0, fy + 2.5, 0.85, 0.34, None, fill=LIGHTB, line=MIDGRAY)
label(s, fx + 2.9, fy + 0.1, 2.6, 0.3, "(b) workflow", size=18, color=GRAY, align=PP_ALIGN.LEFT)
for k in range(4):
    box(s, fx + 3.3, fy + 0.55 + k * 0.9, 1.9, 0.5, None, fill=LIGHTB, line=MIDGRAY)
    if k < 3:
        arrow(s, fx + 4.25, fy + 1.05 + k * 0.9, fx + 4.25, fy + 1.45 + k * 0.9, color=MIDGRAY, weight=1.0)

# ---- Fig 2: search ensembles in hazard space ----
s = fig_slide(2, "search ensembles in hazard space",
    [("b", "shows: each of the three designs' realized search ensembles over the **stationary population's** hazard cloud, with coverage statistics (L2-star discrepancy, MST edge spread)"),
     ("b", "comparison: what each design actually samples; the visual case for random sampling vs hazard coverage"),
     ("b", "role: makes the design tangible before any outcome is shown"),
     ("gap", 2),
     ("c", "processing: hazard image per design draw; discrepancy + spread vs a random design at the same (N, m).")],
    ("Small multiples, one per design, over the SINGLE stationary population's hazard cloud "
     "(there is no shared pool, but all designs share the one population's support). Historic shown "
     "as its own trace marker. Processing: hazard coordinates for every staged search ensemble; "
     "coverage metrics per draw, reported against the random-design expectation."))
fx, fy, fw, fh = mock_frame(s)
panel_names = ["historic", "fixed prob.", "hazard-fill"]
_r.seed(9)
for idx, nm in enumerate(panel_names):
    px = fx + 0.35 + idx * 1.85
    pyy = fy + 1.4
    pw2, ph2 = 1.6, 1.25
    box(s, px, pyy, pw2, ph2, None, fill=WHITE, line=MIDGRAY, line_w=0.75)
    for _ in range(14):
        dot(s, px + 0.07 + _r.random() * (pw2 - 0.2), pyy + 0.07 + _r.random() * (ph2 - 0.2),
            d=0.045, color=MIDGRAY)
    if nm == "historic":
        polyline(s, [(px + 0.25, pyy + 0.8), (px + 0.6, pyy + 0.45), (px + 0.95, pyy + 0.75), (px + 1.3, pyy + 0.5)], ACCENT, 1.75)
    elif nm == "fixed prob.":
        for fx2, fy2 in [(0.4, 0.5), (0.55, 0.3), (0.35, 0.75), (0.6, 0.65), (0.5, 0.45), (0.7, 0.5)]:
            dot(s, px + fx2 * (pw2 - 0.2), pyy + fy2 * (ph2 - 0.2), d=0.08, color=ACCENT)
    else:  # the hazard-filling panel
        for gx2 in range(3):
            for gy2 in range(3):
                dot(s, px + 0.2 + gx2 * (pw2 - 0.55) / 2, pyy + 0.15 + gy2 * (ph2 - 0.45) / 2,
                    d=0.08, color=RUST)
    label(s, px - 0.1, pyy + ph2 + 0.02, pw2 + 0.2, 0.3, nm, size=18, color=GRAY)

# ---- Fig 3: robustness by design ----
s = fig_slide(3, "held-out robustness by design (RQ1 headline)",
    [("b", "shows: distributions of re-evaluated **multivariate satisficing** per design, over K draws × S seeds; univariate panels alongside"),
     ("b", "comparison: the RQ1 verdict, with draw-level spread separating design effect from sampling noise"),
     ("b", "co-reported: the **raw performance distribution** of the same objectives (the degeneracy check)"),
     ("gap", 2),
     ("c", "style cue: Cohen et al. 2021 Fig. 5; Gold et al. 2023 Fig. 5 for the co-reported distributions.")],
    ("Boxplots (or CDFs) of satisficing robustness per design; hazard-filling highlighted. "
     "Processing: per-draw robustness table (design x draw x seed x objective) from the persisted "
     "re-eval cube; mixed-effects variance components with the draw as the unit."))
fx, fy, fw, fh = mock_frame(s)
arrow(s, fx + 0.55, fy + fh - 0.55, fx + 0.55, fy + 0.3, color=MIDGRAY, weight=1.2)
line(s, fx + 0.55, fy + fh - 0.55, fx + fw - 0.25, fy + fh - 0.55, color=MIDGRAY, weight=1.2)
ylab = label(s, fx - 0.9, fy + 1.7, 2.2, 0.35, "satisficing", size=18, color=GRAY)
ylab.rotation = 270
short = ["hist", "fixp", "hz-fill"]
heights = [1.4, 1.9, 2.8]
for i, (nm, hgt) in enumerate(zip(short, heights)):
    cx = fx + 1.35 + i * 1.4
    ybase = fy + fh - 0.55
    line(s, cx, ybase - hgt - 0.45, cx, ybase - hgt + 0.75, color=GRAY, weight=1.0)
    bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.14), Inches(ybase - hgt - 0.2), Inches(0.28), Inches(0.5))
    hzf = nm.startswith("hz")
    bb.fill.solid(); bb.fill.fore_color.rgb = (ACCENT if hzf else LIGHTB)
    bb.line.color.rgb = (ACCENT if hzf else GRAY); bb.line.width = Pt(1.0); bb.shadow.inherit = False
    line(s, cx - 0.14, ybase - hgt + 0.05, cx + 0.14, ybase - hgt + 0.05, color=(WHITE if hzf else INK), weight=1.25)
    label(s, cx - 0.36, ybase + 0.05, 0.72, 0.3, nm, size=18, color=GRAY, wrap=False)
label(s, fx + 0.6, fy + 0.12, fw - 0.8, 0.32, "one panel per objective (7 small multiples)", size=18, color=GRAY)

# ---- Fig 4: threshold sweep ----
s = fig_slide(4, "design robustness vs threshold stringency (RQ1)",
    [("b", "shows: each design's held-out robustness across a grid of satisficing criteria, plus Kendall's **τ_b of the design ranking** vs the default threshold"),
     ("b", "comparison: is the **design ranking** threshold-invariant, or does it flip at the conservative end (Quinn et al. 2020)?"),
     ("b", "role: the credibility figure — a single fixed threshold could manufacture or hide the entire result"),
     ("gap", 2),
     ("c", "processing: `threshold_spectrum` re-scores the persisted cube offline; no re-simulation.")],
    ("Panel (a): robustness vs stringency, one line per design; hazard-filling accented. Panel (b): "
     "tau_b of the design ranking vs the default threshold, across the same grid, with the "
     "stability band annotated. Replaces the deleted robustness-vs-regret figure: regret is not "
     "computed anywhere in this study."))
fx, fy, fw, fh = mock_frame(s)
arrow(s, fx + 0.6, fy + 2.9, fx + 0.6, fy + 0.35, color=MIDGRAY, weight=1.2)
arrow(s, fx + 0.6, fy + 2.9, fx + fw - 0.3, fy + 2.9, color=MIDGRAY, weight=1.2)
ylab = label(s, fx - 0.85, fy + 1.3, 2.2, 0.35, "robustness", size=18, color=GRAY)
ylab.rotation = 270
sweep = [
    ([(0.0, 0.55), (0.33, 0.95), (0.66, 1.65), (1.0, 2.35)], MIDGRAY, 1.5),
    ([(0.0, 0.80), (0.33, 1.25), (0.66, 2.05), (1.0, 2.55)], MIDGRAY, 1.5),
    ([(0.0, 1.05), (0.33, 1.55), (0.66, 1.85), (1.0, 1.95)], MIDGRAY, 1.5),
    ([(0.0, 0.70), (0.33, 0.90), (0.66, 1.05), (1.0, 1.20)], RUST, 2.4),
]
for pts_s, col, wgt in sweep:
    for k in range(len(pts_s) - 1):
        line(s, fx + 0.85 + pts_s[k][0] * 3.9, fy + 0.4 + pts_s[k][1],
             fx + 0.85 + pts_s[k + 1][0] * 3.9, fy + 0.4 + pts_s[k + 1][1], color=col, weight=wgt)
label(s, fx + 3.6, fy + 1.35, 1.7, 0.32, "hazard-fill", size=18, color=RUST, align=PP_ALIGN.LEFT)
label(s, fx + 0.7, fy + 2.95, fw - 1.0, 0.32, "criterion stringency →", size=18, color=GRAY)
label(s, fx + 0.15, fy + 3.35, 2.6, 0.3, "(b) τ_b of the design ranking", size=18, color=GRAY, align=PP_ALIGN.LEFT)
line(s, fx + 0.85, fy + 4.15, fx + 4.75, fy + 4.15, color=MIDGRAY, weight=1.0, dash="dash")
polyline(s, [(fx + 0.85, fy + 3.95), (fx + 2.15, fy + 4.05), (fx + 3.45, fy + 4.3), (fx + 4.75, fy + 4.15)], ACCENT, 2.0)

# ---- Fig 5: the controlled contrast ----
s = fig_slide(5, "the controlled contrast (RQ1)",
    [("b", "shows: re-evaluated robustness distributions for the single contrast fixed_probabilistic vs hazard_filling, with historic as a reference marker"),
     ("b", "comparison: the **selection-rule** effect, read directly — same generator, population, N, L, NFE; only the rule differs"),
     ("b", "role: makes the exact-control argument visible rather than asserted"),
     ("gap", 2),
     ("c", "processing: the same per-draw robustness table as Fig. 3, arranged by the contrast.")],
    ("Two stacked robustness distributions: fixed_probabilistic (top) and hazard_filling (bottom), "
     "the vertical contrast annotated 'only the selection RULE differs'. historic drawn as a "
     "reference line. Replaces the deleted 2x2 population-by-selection and budget-interaction "
     "figures: there is one stationary population and one budget condition."))
fx, fy, fw, fh = mock_frame(s)
cellw, cellh = 4.6, 1.55
for ciy in range(2):
    cx0 = fx + 0.65
    cy0 = fy + 0.75 + ciy * (cellh + 0.7)
    hz = (ciy == 1)
    box(s, cx0, cy0, cellw, cellh, None, fill=WHITE,
        line=(ACCENT if hz else MIDGRAY), line_w=(1.75 if hz else 1.0))
    for k in range(5):
        hgt = 0.4 + 0.14 * k + (0.4 if hz else 0.0)
        bx = cx0 + 0.5 + k * 0.78
        bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx), Inches(cy0 + cellh - 0.18 - hgt),
                                Inches(0.34), Inches(hgt))
        bb.fill.solid(); bb.fill.fore_color.rgb = (RUST if hz else LIGHTB)
        bb.line.fill.background(); bb.shadow.inherit = False
    label(s, cx0 + cellw + 0.05, cy0 + cellh / 2 - 0.2, 1.2, 0.4,
          "hazard-fill" if hz else "fixed prob.", size=18, color=(RUST if hz else GRAY), align=PP_ALIGN.LEFT)
arrow(s, fx + 1.5, fy + 2.3, fx + 1.5, fy + 2.75, color=ACCENT, weight=1.9)
label(s, fx + 1.75, fy + 2.32, 3.2, 0.4, "only the selection RULE differs", size=18, color=ACCENT, align=PP_ALIGN.LEFT)
label(s, fx + 0.35, fy + 4.6, fw - 0.7, 0.4,
      "historic shown as a reference line", size=18, color=GRAY)

# ---- Fig 6: policy differences ----
s = fig_slide(6, "do designs find different policies? (RQ1)",
    [("b", "shows: re-evaluated objective tradeoffs as parallel coordinates, policies colored by design; compromise policies highlighted"),
     ("b", "comparison: composition can move **values more than rankings** (McPhail et al. 2020); do designs change the decisions, not just the scores?"),
     ("gap", 2),
     ("c", "style cue: Trindade et al. 2017 Fig. 6 (robustness parallel coordinates with highlighted compromises).")],
    ("Parallel coordinates over the 7 re-evaluated objectives; gray cloud = all policies, colored "
     "bundles by design, thick lines = per-design compromise policies. A companion DV-space "
     "version (decision-variable ranges per design, Trindade Fig. 8 style) may go in supplement. "
     "Processing: re-evaluated objective matrix + archived decision variables per design."))
fx, fy, fw, fh = mock_frame(s)
n_ax = 7
ax_xs = [fx + 0.5 + k * (fw - 1.0) / (n_ax - 1) for k in range(n_ax)]
for k, axx in enumerate(ax_xs):
    line(s, axx, fy + 0.55, axx, fy + fh - 0.75, color=MIDGRAY, weight=1.2)
    label(s, axx - 0.25, fy + fh - 0.68, 0.5, 0.3, "f%d" % (k + 1), size=18, color=GRAY)
_r.seed(14)
for _ in range(7):
    pts6 = [(axx, fy + 0.7 + _r.random() * (fh - 1.6)) for axx in ax_xs]
    polyline(s, pts6, MIDGRAY, 0.9)
for col, wgt in [(ACCENT, 2.0), (RUST, 2.5)]:
    pts6 = [(axx, fy + 0.7 + _r.random() * (fh - 1.6)) for axx in ax_xs]
    polyline(s, pts6, col, wgt)
label(s, fx + 0.5, fy + 0.12, fw - 0.8, 0.35, "colored: compromise policies by design", size=18, color=GRAY)

# ---- Fig 7: mechanism test ----
s = fig_slide(7, "mechanism test: failure vs coverage deficit (RQ1)",
    [("b", "shows: per design, the held-out failure rate as a function of **hazard-space coverage deficit**, against a **random-ensemble null**"),
     ("b", "comparison: the falsifiable prediction — hazard-filling's policies should fail in the region its own design under-covered"),
     ("b", "caveat drawn on the figure: raw AUC is inflated by geometry, so the verdict is **AUC minus the null**, never raw AUC"),
     ("gap", 2),
     ("c", "processing: boosted-tree scenario discovery in hazard space (Gold et al. 2022, 2023 configuration); null from random ensembles at the same (N, m).")],
    ("Panel (a): failure rate vs coverage deficit, one curve per design, with the random-ensemble "
     "null shaded. Panel (b): AUC minus null, per design, with a CI. Replaces the old "
     "coverage-vs-robustness scatter, which was correlational; this is the mechanism test."))
fx, fy, fw, fh = mock_frame(s)
arrow(s, fx + 0.6, fy + fh - 0.6, fx + 0.6, fy + 0.35, color=MIDGRAY, weight=1.2)
arrow(s, fx + 0.6, fy + fh - 0.6, fx + fw - 0.3, fy + fh - 0.6, color=MIDGRAY, weight=1.2)
ylab = label(s, fx - 0.9, fy + 1.75, 2.2, 0.35, "failure rate", size=18, color=GRAY)
ylab.rotation = 270
label(s, fx + 0.65, fy + fh - 0.48, fw - 1.0, 0.32, "hazard-space coverage deficit →", size=18, color=GRAY)
null_pts = [(0.9, 3.35), (2.0, 3.05), (3.1, 2.80), (4.6, 2.55)]
for k in range(len(null_pts) - 1):
    ln2 = line(s, fx + null_pts[k][0], fy + null_pts[k][1],
               fx + null_pts[k + 1][0], fy + null_pts[k + 1][1], color=MIDGRAY, weight=1.4)
    ln_el = ln2.line._get_or_add_ln(); d_el = ln_el.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln_el.append(d_el)
label(s, fx + 3.3, fy + 2.15, 2.1, 0.3, "random null", size=18, color=MIDGRAY, align=PP_ALIGN.LEFT)
polyline(s, [(fx + 0.9, fy + 3.55), (fx + 2.0, fy + 2.75), (fx + 3.1, fy + 1.85), (fx + 4.6, fy + 1.05)], RUST, 2.3)
label(s, fx + 3.05, fy + 0.65, 2.4, 0.3, "hazard-fill", size=18, color=RUST, align=PP_ALIGN.LEFT)
_r.seed(4)
for _ in range(12):
    xx = 0.95 + _r.random() * 3.6
    yy = fy + 3.6 - (xx - 0.9) * 0.6 - _r.random() * 0.6
    dot(s, fx + xx, yy, d=0.08, color=MIDGRAY)

# ---- Fig 8: RQ2 vs current FFMP ----
s = fig_slide(8, "re-optimized vs current FFMP (RQ2)",
    [("b", "shows: (a) objective-by-objective change vs the current FFMP baseline on the test ensemble; (b) storage / Montague flow through a drought-of-record-scale event"),
     ("b", "comparison: Pareto-approximate policies vs the status-quo baseline — the same fixed external reference used by the improvement-over-status-quo metric"),
     ("gap", 2),
     ("c", "processing: baseline FFMP re-evaluation (workflow step 05); event extraction from test scenarios; policy simulation traces.")],
    ("Panel (a): horizontal bars of change vs baseline per objective (improvement right). Panel "
     "(b): time series through a severe drought, current FFMP vs a compromise re-optimized policy, "
     "drought window shaded. Processing: current-FFMP run on E_test; trace storage for selected "
     "policies and events."))
fx, fy, fw, fh = mock_frame(s)
label(s, fx + 0.15, fy + 0.1, 2.5, 0.3, "(a) change vs FFMP", size=18, color=GRAY, align=PP_ALIGN.LEFT)
zx = fx + 1.3
line(s, zx, fy + 0.55, zx, fy + fh - 0.4, color=GRAY, weight=1.2)
bars = [0.9, 0.6, 1.1, 0.5, -0.35, 0.75, 0.85]
for k, bl in enumerate(bars):
    yy = fy + 0.7 + k * 0.5
    bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(zx if bl >= 0 else zx + bl * 1.0), Inches(yy),
                            Inches(abs(bl) * 1.0), Inches(0.24))
    bb.fill.solid(); bb.fill.fore_color.rgb = (ACCENT if bl >= 0 else RUST)
    bb.line.fill.background(); bb.shadow.inherit = False
    label(s, fx + 0.15, yy - 0.04, 0.55, 0.3, "f%d" % (k + 1), size=18, color=GRAY, align=PP_ALIGN.LEFT)
label(s, fx + 3.0, fy + 0.1, 2.6, 0.3, "(b) drought event", size=18, color=GRAY, align=PP_ALIGN.LEFT)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(fx + 3.7), Inches(fy + 0.6), Inches(1.1), Inches(3.1))
band.fill.solid(); band.fill.fore_color.rgb = LIGHTB; band.line.fill.background(); band.shadow.inherit = False
arrow(s, fx + 3.1, fy + 3.7, fx + 3.1, fy + 0.55, color=MIDGRAY, weight=1.2)
line(s, fx + 3.1, fy + 3.7, fx + 5.45, fy + 3.7, color=MIDGRAY, weight=1.2)
polyline(s, [(fx + 3.1, fy + 1.1), (fx + 3.7, fy + 1.35), (fx + 4.1, fy + 2.9), (fx + 4.6, fy + 3.3), (fx + 5.0, fy + 2.1), (fx + 5.4, fy + 1.4)], RUST, 1.75)
polyline(s, [(fx + 3.1, fy + 1.0), (fx + 3.7, fy + 1.2), (fx + 4.1, fy + 2.1), (fx + 4.6, fy + 2.45), (fx + 5.0, fy + 1.6), (fx + 5.4, fy + 1.15)], ACCENT, 1.75)
label(s, fx + 3.15, fy + 3.85, 2.4, 0.3, "storage: current (rust) vs re-optimized", size=18, color=GRAY, align=PP_ALIGN.LEFT)

# ---- Fig 9: RQ3 rule resolution ----
s = fig_slide(9, "rule resolution, ffmp_N (RQ3)",
    [("b", "shows: re-evaluated robustness vs FFMP storage-zone resolution N, including the 24-parameter standard"),
     ("b", "comparison: is added rule resolution worth its larger search space?"),
     ("gap", 2),
     ("c", "processing: the ffmp_N campaign runs; same re-evaluation post-processing as RQ1. Per-design (never pooled) reference sets if a search metric is shown.")],
    ("Dot-line plot: x = rule resolution (8, 10, 12, standard), y = re-evaluated satisficing. Any "
     "MOEA search metric (hypervolume) is a WITHIN-design convergence diagnostic only, per "
     "Zatarain Salazar 2017; it is not a cross-design yardstick. Processing: ffmp_N sweep under the "
     "chosen design(s), re-evaluated on E_test."))
fx, fy, fw, fh = mock_frame(s)
arrow(s, fx + 0.55, fy + fh - 0.55, fx + 0.55, fy + 0.3, color=MIDGRAY, weight=1.2)
line(s, fx + 0.55, fy + fh - 0.55, fx + fw - 0.25, fy + fh - 0.55, color=MIDGRAY, weight=1.2)
ylab = label(s, fx - 0.9, fy + 1.7, 2.2, 0.35, "satisficing", size=18, color=GRAY)
ylab.rotation = 270
ticks = ["8", "10", "12", "std"]
txs = [fx + 1.2 + k * 1.15 for k in range(4)]
for tx, tl in zip(txs, ticks):
    label(s, tx - 0.3, fy + fh - 0.45, 0.6, 0.3, tl, size=18, color=GRAY)
ys1 = [2.9, 2.4, 2.2, 2.3]
polyline(s, list(zip(txs, [fy + yv for yv in ys1])), ACCENT, 2.0)
for tx, yv in zip(txs, ys1):
    dot(s, tx - 0.05, fy + yv - 0.05, d=0.11, color=ACCENT)
label(s, fx + 0.6, fy + 0.12, fw - 0.8, 0.32, "held-out satisficing, per rule resolution", size=18, color=GRAY)
label(s, fx + 0.55, fy + fh - 0.15, fw - 0.8, 0.32, "FFMP storage-zone resolution", size=18, color=GRAY)

OUT = r"c:\Users\tjame\Desktop\Research\DRB\Pywr-DRB\NYCOptimization\docs\NYCOptimization_project_proposal_and_motivation_CLAUDE.pptx"
prs.save(OUT)
print("saved", OUT, "n_slides:", len(prs.slides._sldIdLst))

# -*- coding: utf-8 -*-
"""Build NYCOpt_Objective_Definitions_CLAUDE.pptx — committee proposal draft.

Style: minimalist, Cambria titles / Calibri body, all text >= 18 pt,
simple native-shape diagrams, muted palette.
"""
import random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------- palette ----------
INK    = RGBColor(0x21, 0x25, 0x29)   # near-black text
ACCENT = RGBColor(0x2E, 0x5E, 0x8C)   # steel blue
GRAY   = RGBColor(0x5F, 0x6B, 0x7A)   # citations / labels
LIGHTB = RGBColor(0xE9, 0xEF, 0xF5)   # light fill
RUST   = RGBColor(0xA4, 0x59, 0x29)   # sparing highlight (proposed design)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MIDGRAY= RGBColor(0x9A, 0xA4, 0xAF)

TITLE_FONT = "Cambria"
BODY_FONT  = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------

def add_slide(notes=None):
    s = prs.slides.add_slide(BLANK)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _apply_runs(p, text, size, color, font=BODY_FONT, italic=False, base_bold=False):
    """Add runs to paragraph p, parsing **bold** inline markup."""
    parts = text.split("**")
    for i, part in enumerate(parts):
        if not part:
            continue
        r = p.add_run()
        r.text = part
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.italic = italic
        r.font.bold = base_bold or (i % 2 == 1)


def add_title(s, text, size=28):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12.13), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _apply_runs(p, text, size, INK, font=TITLE_FONT)
    # accent rule under title
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.63), Inches(1.18), Inches(1.9), Pt(2.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return tb


def add_body(s, items, left=0.6, top=1.5, width=12.13, height=None,
             size=20, space_after=10, line_spacing=1.04):
    """items: list of (kind, text) where kind in
    'b' bullet, 'b2' sub-bullet, 'p' plain, 'c' citation, 'gap' spacer,
    'h' bold lead-in plain line."""
    if height is None:
        height = 7.25 - top
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for kind, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if kind == "gap":
            p.space_after = Pt(text if isinstance(text, (int, float)) else 6)
            continue
        if kind == "b":
            _apply_runs(p, u"•  " + text, size, INK)
        elif kind == "b2":
            p.level = 1
            _apply_runs(p, u"–  " + text, size - 1 if size > 18.5 else 18, INK)
        elif kind == "p":
            _apply_runs(p, text, size, INK)
        elif kind == "h":
            _apply_runs(p, text, size, INK, base_bold=True)
        elif kind == "c":
            _apply_runs(p, text, 18, GRAY, italic=True)
        elif kind == "f":  # formula-style line, Cambria accent-dark
            _apply_runs(p, text, size, INK, font=TITLE_FONT)
    return tb


def box(s, x, y, w, h, text=None, fill=LIGHTB, line=GRAY, size=18, bold=False,
        color=INK, shape=MSO_SHAPE.RECTANGLE, align=PP_ALIGN.CENTER, line_w=1.0,
        font=BODY_FONT):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if text is not None:
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        _apply_runs(p, text, size, color, font=font, base_bold=bold)
    return sh


def label(s, x, y, w, h, text, size=18, color=GRAY, align=PP_ALIGN.CENTER,
          italic=False, bold=False, font=BODY_FONT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _apply_runs(p, text, size, color, font=font, italic=italic, base_bold=bold)
    return tb


def arrow(s, x1, y1, x2, y2, color=GRAY, weight=1.5):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def dot(s, x, y, d=0.09, color=ACCENT):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def axes_panel(s, x, y, w, h, xlabel, ylabel):
    """L-shaped axes with labels below/left."""
    arrow(s, x, y + h, x + w, y + h, color=MIDGRAY, weight=1.2)   # x axis
    arrow(s, x, y + h, x, y, color=MIDGRAY, weight=1.2)           # y axis
    label(s, x, y + h + 0.03, w, 0.32, xlabel, size=18, color=GRAY)
    lb = label(s, x - 2.05, y + h / 2 - 0.17, 2.0, 0.34, ylabel, size=18, color=GRAY, align=PP_ALIGN.RIGHT)
    return lb


NO_STYLE_TABLE = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # No Style, Table Grid

def add_table(s, rows, col_widths, left=0.6, top=1.5, row_h=0.5,
              size=18, header_fill=LIGHTB, highlight_row=None, highlight_color=None):
    n_rows, n_cols = len(rows), len(rows[0])
    total_w = sum(col_widths)
    gfx = s.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                             Inches(total_w), Inches(row_h * n_rows))
    tbl = gfx.table
    # neutral table style (thin grid, no banding)
    tblPr = tbl._tbl.tblPr
    tblPr.set('bandRow', '0'); tblPr.set('firstRow', '1')
    for el in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(el)
    style_el = tblPr.makeelement(qn('a:tableStyleId'), {})
    style_el.text = NO_STYLE_TABLE
    tblPr.append(style_el)
    for j, wdt in enumerate(col_widths):
        tbl.columns[j].width = Inches(wdt)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(row_h)
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            elif highlight_row is not None and i == highlight_row:
                cell.fill.solid(); cell.fill.fore_color.rgb = (highlight_color or LIGHTB)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            color = INK
            _apply_runs(p, cell_text, size, color, base_bold=(i == 0))
    return gfx


def divider(s, text_lines):
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = LIGHTB
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, bold) in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _apply_runs(p, txt, 32, INK, font=TITLE_FONT, base_bold=bold)


# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = add_slide(notes=(
    "Working title - two alternates to consider: 'Designing the Scenarios We Optimize Over' or "
    "'Hazard-Space Scenario Design for Many-Objective Reservoir Policy Search'. "
    "Framing: this is a dissertation-paper proposal; no results yet. The deliverable today is the "
    "scope, the gap, and the planned experiment."))
label(s, 0.9, 2.15, 11.5, 1.6,
      "Hazard-Informed Scenario Design for\nMany-Objective Reservoir Policy Search",
      size=36, color=INK, align=PP_ALIGN.LEFT, font=TITLE_FONT)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(3.62), Inches(2.6), Pt(3))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
label(s, 0.9, 3.85, 11.5, 0.5,
      "Re-optimizing NYC reservoir operations in the Delaware River Basin under hydrologic uncertainty",
      size=22, color=GRAY, align=PP_ALIGN.LEFT)
label(s, 0.9, 5.35, 11.5, 1.2,
      "Research proposal  |  Trevor Amestoy\nReed Research Group, Cornell University  |  Summer 2026",
      size=20, color=INK, align=PP_ALIGN.LEFT)

# =====================================================================
# SLIDE 2 — Overview
# =====================================================================
s = add_slide(notes=(
    "Roadmap. Part 1 motivates the problem and states the gap; Part 2 is the proposal proper "
    "(hypothesis, RQs, contributions); Part 3 is the planned experiment; close with status and the "
    "specific decisions where committee input is most valuable."))
add_title(s, "Overview")
add_body(s, [
    ("h", "1.  Motivation & gap"),
    ("b", "How the streamflow scenarios used during policy search are chosen today, and why that choice matters"),
    ("gap", 4),
    ("h", "2.  Proposal"),
    ("b", "Hypothesis, research questions, and planned contributions"),
    ("gap", 4),
    ("h", "3.  Planned experiment"),
    ("b", "Testbed, scenario designs, objectives, comparison controls, evaluation"),
    ("gap", 4),
    ("h", "4.  Status, open decisions & requested feedback"),
    ("gap", 4),
    ("c", "Supplemental: objective formulations, ensemble aggregation, forcing parameterization, HPC scaling."),
], top=1.55, size=20)

# =====================================================================
# SLIDE 3 — Testbed: NYC reservoirs in the DRB
# =====================================================================
s = add_slide(notes=(
    "Orient the committee to the system before any methodology. Key facts: three NYC Delaware "
    "reservoirs; the 1954 Supreme Court Decree fixes the diversion cap and the downstream flow "
    "targets; day-to-day operations follow the FFMP. All simulation is in Pywr-DRB. The schematic "
    "is deliberately minimal - reservoirs, diversion, two Decree flow targets."))
add_title(s, "Testbed — NYC reservoirs in the Delaware River Basin")
add_body(s, [
    ("b", "Cannonsville, Pepacton & Neversink supply roughly half of NYC's water and control flows into the upper Delaware mainstem."),
    ("b", "1954 U.S. Supreme Court Decree: NYC diversion capped at 800 MGD; minimum-flow obligations at Montague (1,750 cfs) and Trenton protect New Jersey and Philadelphia (salinity repulsion)."),
    ("b", "Operations follow the Flexible Flow Management Program (FFMP): releases set by storage zone, with drought step-downs."),
    ("b", "Simulated with Pywr-DRB, an open-source daily water-availability model of the full basin."),
    ("gap", 2),
    ("c", "Hamilton et al. 2024; Kolesar & Serio 2011."),
], left=0.6, top=1.5, width=6.3, size=20)
# schematic (right column)
rx = 8.55
box(s, rx + 0.0, 1.7, 1.5, 0.62, "Cannonsville", size=18)
box(s, rx + 1.62, 1.7, 1.5, 0.62, "Pepacton", size=18)
box(s, rx + 3.24, 1.7, 1.5, 0.62, "Neversink", size=18)
arrow(s, rx + 0.75, 2.32, rx + 2.2, 3.15)
arrow(s, rx + 2.37, 2.32, rx + 2.37, 3.15)
arrow(s, rx + 3.99, 2.32, rx + 2.55, 3.15)
box(s, rx + 1.45, 3.2, 1.9, 0.6, "Delaware mainstem", fill=None, line=MIDGRAY, size=18)
arrow(s, rx + 2.4, 3.8, rx + 2.4, 4.35)
box(s, rx + 1.2, 4.4, 2.3, 0.72, "Montague target\n1,750 cfs", size=18)
arrow(s, rx + 2.35, 5.12, rx + 2.35, 5.65)
box(s, rx + 1.2, 5.7, 2.3, 0.72, "Trenton target\n3,000 cfs", size=18)
# NYC diversion arrow from reservoir row to left label
box(s, 7.05, 2.75, 1.6, 0.85, "NYC diversion\n≤ 800 MGD", fill=WHITE, line=ACCENT, size=18)
arrow(s, rx + 0.5, 2.35, 8.3, 2.8)

# =====================================================================
# SLIDE 4 — Policy search under uncertainty
# =====================================================================
s = add_slide(notes=(
    "Set up the general MORDM-style workflow so the object of study is unambiguous: the evaluation "
    "ensemble inside the search loop. Emphasize the asymmetry: the optimizer only ever 'sees' "
    "hydrology through this ensemble. Re-evaluation happens after search, on broader conditions."))
add_title(s, "Policy search under uncertainty — where scenarios enter")
add_body(s, [
    ("b", "Many-objective search: the MOEA proposes candidate operating policies; each is simulated over an **evaluation ensemble** of streamflow scenarios to compute its objectives."),
    ("b", "The evaluation ensemble is the optimizer's entire view of hydrology — policies are only rewarded for handling conditions that the ensemble actually contains."),
    ("b", "After search, Pareto-approximate policies are re-evaluated on broader, out-of-sample conditions to measure robustness."),
    ("b", "**This proposal asks: how should the search-time evaluation ensemble be designed?**"),
    ("gap", 2),
    ("c", "Kasprzyk et al. 2013; Herman et al. 2015, 2020."),
], top=1.5, size=20)
dy = 5.3
box(s, 0.7, dy, 2.5, 0.85, "MM-Borg MOEA\ncandidate policy", size=18)
arrow(s, 3.2, dy + 0.42, 3.8, dy + 0.42)
box(s, 3.8, dy, 3.4, 0.85, "Pywr-DRB simulation over\n**evaluation ensemble**", size=18, line=ACCENT)
arrow(s, 7.2, dy + 0.42, 7.8, dy + 0.42)
box(s, 7.8, dy, 2.3, 0.85, "objective values", size=18)
arrow(s, 10.1, dy + 0.42, 10.7, dy + 0.42)
box(s, 10.7, dy, 2.1, 0.85, "re-evaluation\n(after search)", fill=None, line=MIDGRAY, size=18)
# feedback arrow (objectives -> MOEA)
arrow(s, 8.9, dy + 1.05, 1.95, dy + 1.05, color=MIDGRAY, weight=1.2)
label(s, 4.0, dy + 1.12, 3.6, 0.35, "search loop", size=18, color=GRAY)

# =====================================================================
# SLIDE 5 — Motivation 1: beyond the historical record
# =====================================================================
s = add_slide(notes=(
    "Two linked points: (1) the historical record undersamples extremes, so synthetic ensembles are "
    "standard; (2) but naive large ensembles spend most of their simulated years in unremarkable "
    "conditions - the 84-93% figure is from our current DRB stochastic ensembles. The search signal "
    "comes from the rare stress events."))
add_title(s, "Motivation — ensembles are necessary, but inefficient")
add_body(s, [
    ("b", "The historical record is one short draw from a variable climate: it undersamples droughts and floods beyond experience, so search over the record alone overfits to one hydrologic history."),
    ("b", "Synthetic streamflow ensembles are the standard remedy — they expand the range of variability a policy is tested against."),
    ("b", "But large stochastic ensembles are dominated by unremarkable hydrology: **84–93% of years in our current DRB ensembles fall in normal conditions** — simulated years that add cost, not search signal."),
    ("b", "Rare, severe events drive the consequential outcomes; the optimizer learns from the few scenarios that contain them."),
    ("gap", 2),
    ("c", "Kirsch et al. 2013; Herman et al. 2016; Cohen et al. 2021; Zaniolo et al. 2023."),
], top=1.5, size=20)

# =====================================================================
# SLIDE 6 — Motivation 2: input-space sampling is redundant
# =====================================================================
s = add_slide(notes=(
    "The core inefficiency argument. Prevailing DMDU practice samples the input space (generator / "
    "climate-change parameters). But the mapping from parameters to realized hydrologic stress is "
    "many-to-one and noisy: distinct parameter vectors yield overlapping hazard conditions. Our "
    "preliminary DRB diagnostic found an input-vs-hazard coverage gap (~1.3x discrepancy) that "
    "survives even enriched input sampling. Diagram: spread points in parameter space collapse to "
    "overlapping regions in hazard space."))
add_title(s, "Motivation — input-space sampling gives redundant scenarios")
add_body(s, [
    ("b", "Prevailing deep-uncertainty practice samples the **input space**: Latin hypercube designs over generator / climate parameters, one 'state of the world' per sample."),
    ("b", "Distinct parameter sets frequently produce hydrologically redundant realizations — uniform coverage of inputs does not give uniform coverage of the hazard conditions that stress the system."),
    ("b", "Preliminary DRB diagnostic: hazard-space coverage gap of input-space designs persists (~1.3× discrepancy) even after enriching the input sample."),
    ("gap", 2),
    ("c", "Quinn et al. 2018; Bartholomew & Kwakkel 2020; Guo et al. 2018; Quinn et al. 2020."),
], top=1.5, size=20)
# diagram: input space dots -> hazard space overlapping ovals
py = 5.15; ph = 1.55; pw = 2.6
axes_panel(s, 2.6, py, pw, ph, "generator parameter 1", "parameter 2")
for (fx, fy) in [(0.15, 0.2), (0.5, 0.75), (0.8, 0.35), (0.3, 0.55), (0.65, 0.12), (0.9, 0.8), (0.12, 0.85)]:
    dot(s, 2.6 + fx * (pw - 0.25) + 0.06, py + fy * (ph - 0.25) + 0.06, d=0.11, color=ACCENT)
arrow(s, 5.75, py + ph / 2, 6.95, py + ph / 2, color=GRAY, weight=1.8)
label(s, 5.3, py + ph / 2 - 0.82, 2.1, 0.7, "generate\nstreamflow", size=18, color=GRAY)
axes_panel(s, 8.9, py, pw, ph, "hazard metric 1", "metric 2")
for (ox, oy, ow, oh) in [(0.55, 0.45, 1.3, 0.72), (0.75, 0.3, 1.3, 0.72), (0.45, 0.62, 1.3, 0.72)]:
    e = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.9 + ox), Inches(py + oy - 0.15), Inches(ow), Inches(oh))
    e.fill.solid(); e.fill.fore_color.rgb = LIGHTB
    e.line.color.rgb = ACCENT; e.line.width = Pt(1.0); e.shadow.inherit = False
label(s, 8.75, py - 0.42, 3.2, 0.36, "overlapping = redundant", size=18, color=RUST, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 7 — Terminology: the three spaces
# =====================================================================
s = add_slide(notes=(
    "This vocabulary is load-bearing for the whole talk. Input space = what defines generation. "
    "Hazard space = metrics measured on each realized sequence BEFORE any system simulation "
    "(drought intensity/duration/severity, low/high-flow indices). Outcome space = simulated "
    "objectives. Because hazard coordinates exist before any policy is evaluated, subsampling in "
    "hazard space is a pre-optimization design step - simulation-free."))
add_title(s, "Terminology — three spaces")
dy = 1.75
box(s, 0.7, dy, 3.5, 1.0, "**Input space**\ngenerator & climate-forcing parameters θ", size=18)
arrow(s, 4.2, dy + 0.5, 4.95, dy + 0.5)
label(s, 3.68, dy + 1.05, 1.8, 0.35, "generate", size=18, color=GRAY)
box(s, 4.95, dy, 3.5, 1.0, "**Hazard space**\nhazard metrics measured on each realized sequence", size=18, line=ACCENT, line_w=1.75)
arrow(s, 8.45, dy + 0.5, 9.2, dy + 0.5)
label(s, 7.93, dy + 1.05, 1.8, 0.35, "simulate", size=18, color=GRAY)
box(s, 9.2, dy, 3.5, 1.0, "**Outcome space**\nobjective values of a policy under a scenario", size=18)
add_body(s, [
    ("b", "**Input space** — factors that define scenario generation (generator parameters, monthly climate change factors). A sampled point is a 'state of the world' in the MORDM literature."),
    ("b", "**Hazard space** — drought event metrics (SSI intensity / duration / severity), low-flow and high-flow indices, computed on each streamflow sequence **before any system simulation**."),
    ("b", "**Outcome space** — simulated performance; reserved strictly for objectives."),
    ("b", "Hazard coordinates exist before any policy is evaluated → designing the ensemble in hazard space is a **simulation-free, pre-optimization step**."),
    ("gap", 2),
    ("c", "Yevjevich 1967; Vicente-Serrano et al. 2012; Richter et al. 1996; Kasprzyk et al. 2013."),
], top=3.35, size=20)

# =====================================================================
# SLIDE 8 — The gap
# =====================================================================
s = add_slide(notes=(
    "The honest gap statement. Not 'nobody asked whether composition matters' - Cohen et al. 2021 "
    "did, affirmatively, and we cite it as motivation. Each adjacent literature stops short: Cohen "
    "(problem-driven regret selection, 97 GCM traces, split-half testing); Bonham (space-filling but "
    "post-hoc ranking); Zaniolo (hazard control at generation, discrete types); Zatarain Salazar "
    "(1-D magnitude stratification). The delta: scalable, simulation-free, coverage-designed "
    "construction of the SEARCH ensemble, tested under genuinely held-out deep-uncertainty "
    "re-evaluation."))
add_title(s, "The gap — adjacent literatures stop short")
add_body(s, [
    ("b", "**Search-set composition matters** — training-scenario properties drive out-of-sample policy robustness (Cohen et al. 2021) — but their selection is problem-driven (a perfect-foresight optimization per scenario) over 97 GCM traces."),
    ("b", "**Space-filling subset selection** entered water resources for **post-hoc robustness ranking**, not the search ensemble (Bonham et al. 2024)."),
    ("b", "**Hazard-targeted ensembles** are generated to a few discrete drought types — control at generation, not selection from realized sequences (Zaniolo et al. 2023; Borgomeo et al. 2015)."),
    ("b", "Nearest search-phase stratification is **one-dimensional** (flow magnitude; Zatarain Salazar et al. 2017)."),
    ("gap", 4),
    ("h", "No published study designs the MOEA evaluation ensemble as a space-filling sample of a multi-dimensional hazard space and tests it under a genuinely held-out, deeply uncertain re-evaluation."),
], top=1.5, size=20)

# =====================================================================
# SLIDE 9 — Proposed idea: hazard-filling
# =====================================================================
s = add_slide(notes=(
    "The idea in one picture, and the argument behind it. Hazard coordinates are EMERGENT "
    "properties of a realized flow sequence -- no generator can be asked to produce a realization "
    "at a prescribed drought severity. Forcing parameters theta ARE a knob. So input-space designs "
    "GENERATE TO their design points (LHS alone, nothing to snap to), while hazard-space designs "
    "must SELECT FROM a finite candidate pool (LHS anchors + nearest-neighbour snap). That "
    "asymmetry is intrinsic, not an implementation shortcut. Each hazard-filling design generates "
    "and owns its candidate pool; nothing is shared between designs. Be explicit about the "
    "probability distortion -- deliberate, and compared only on the held-out re-evaluation."))
add_title(s, "Proposed design — hazard-filling evaluation ensembles")
add_body(s, [
    ("b", "**θ is a knob on the generator; hazard is not** — it is emergent, measured on a realized sequence. So input-space designs **generate to** their design points, while hazard-space designs must **select from** a pool."),
    ("b", "Each hazard-filling design generates and owns a **candidate pool** of ~10⁵–10⁶ short (10-yr) i.i.d. sequences, and computes each sequence's hazard coordinates."),
    ("b", "**Select** the evaluation ensemble by **Latin-hypercube anchors in hazard space, snapped to the nearest unused pool member** — deterministic, no tuning. The snap is intrinsic, not an approximation of something better."),
    ("b", "Simulation-free selection → scalable to the full pool, reusable across formulations."),
    ("b", "Deliberately distorts scenario probabilities toward severe corners; the held-out re-evaluation is the probability-faithful corrective."),
], top=1.5, size=20, space_after=7)
# diagram: dense cloud -> subsample
py = 5.0; ph = 1.7; pw = 2.6
axes_panel(s, 2.6, py, pw, ph, "hazard metric 1", "metric 2")
random.seed(21)
for _ in range(70):
    fx, fy = random.random(), random.random()
    dot(s, 2.6 + 0.1 + fx * (pw - 0.35), py + 0.08 + fy * (ph - 0.35), d=0.055, color=MIDGRAY)
label(s, 2.0, py - 0.45, 3.8, 0.36, "candidate pool (own, i.i.d., ~10⁶)", size=18, color=GRAY)
arrow(s, 5.8, py + ph / 2, 7.0, py + ph / 2, color=GRAY, weight=1.8)
label(s, 5.15, py + ph / 2 - 0.9, 2.8, 0.7, "LHS anchors +\nnearest-neighbour snap", size=18, color=GRAY)
axes_panel(s, 9.0, py, pw, ph, "hazard metric 1", "metric 2")
for fx, fy in [(0.1, 0.15), (0.5, 0.1), (0.9, 0.15), (0.1, 0.5), (0.5, 0.5), (0.88, 0.52),
               (0.12, 0.85), (0.5, 0.88), (0.9, 0.85), (0.3, 0.3), (0.7, 0.32), (0.3, 0.7), (0.7, 0.7)]:
    dot(s, 9.0 + 0.1 + fx * (pw - 0.35), py + 0.08 + fy * (ph - 0.35), d=0.1, color=ACCENT)
label(s, 8.6, py - 0.45, 3.4, 0.36, "evaluation ensemble (N ≈ 10²)", size=18, color=GRAY)

# =====================================================================
# SLIDE 10 — Divider: Proposal
# =====================================================================
s = add_slide(notes="Section divider.")
divider(s, [("Proposal:", False), ("Hypothesis, Research Questions & Contributions", True)])

# =====================================================================
# SLIDE 11 — Hypothesis & contributions
# =====================================================================
s = add_slide(notes=(
    "State the falsifiable hypothesis, then the contributions contingent on results. Note contribution "
    "1 stands even under a null result - the controlled comparison itself is new. Frame relative to "
    "Cohen 2021: they established composition matters; we contribute the scalable, simulation-free, "
    "coverage-designed construction and the held-out deep-uncertainty test."))
add_title(s, "Hypothesis & planned contributions")
hb = box(s, 0.6, 1.55, 12.15, 1.5, None, fill=LIGHTB, line=ACCENT, line_w=1.25)
tf = hb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
_apply_runs(p, "**Hypothesis.**  Evaluation ensembles that uniformly cover hazard space yield Pareto-approximate policies that are more robust under held-out re-evaluation than probabilistic or input-space designs, at equal simulation budget.", 20, INK)
add_body(s, [
    ("h", "Planned contributions"),
    ("b", "**First controlled comparison of scenario designs** for MOEA evaluation ensembles: six designs spanning the literature's main families, at matched computational budget, judged on one common held-out test ensemble."),
    ("b", "**A scalable, simulation-free, coverage-designed construction** of the search ensemble (hazard-filling), with its probability distortion made explicit and its effect measured."),
    ("b", "Evidence on whether **re-optimizing the FFMP rule structure** can improve NYC and basin outcomes, including a variable-resolution variant with more storage zones."),
    ("gap", 2),
    ("c", "Positioning: Cohen et al. 2021 showed composition matters; the contribution here is the construction and its test."),
], top=3.35, size=20)

# =====================================================================
# SLIDE 12 — Research questions
# =====================================================================
s = add_slide(notes=(
    "RQ1 is the core methodological question and drives the experimental design; RQ2 and RQ3 are the "
    "domain questions the same campaign answers along the way."))
add_title(s, "Research questions")
add_body(s, [
    ("h", "RQ1 (core).  Does the composition of the search-time scenario ensemble change the robustness of the resulting Pareto-approximate policies under held-out re-evaluation?"),
    ("gap", 8),
    ("h", "RQ2.  Can re-optimizing the FFMP parameters improve NYC and basin outcomes?"),
    ("b", "Supply reliability, Montague / Trenton Decree targets, downstream flooding, storage resilience."),
    ("gap", 8),
    ("h", "RQ3.  Does a variable-resolution FFMP structure (more storage zones) improve performance or robustness?"),
    ("b", "FFMP with N ∈ {8, 10, 12} storage-zone layers vs the standard 24-variable structure."),
], top=1.7, size=20)

# =====================================================================
# SLIDE 13 — Divider: Planned experiment
# =====================================================================
s = add_slide(notes="Section divider.")
divider(s, [("Planned Experiment", True)])

# =====================================================================
# SLIDE 14 — Experiment at a glance (pipeline)
# =====================================================================
s = add_slide(notes=(
    "The whole experiment in one diagram. Top row: two populations (stationary Kirsch-Nowak at "
    "theta_hist; DU-forced theta from the CMIP6 hypercube). Bottom row: the comparison (each design "
    "builds its OWN ensemble -> independent searches at identical N, L and NFE -> common held-out "
    "re-evaluation). KEY CONTROL: there is no shared pool. Every design generates its own "
    "realizations from its own namespaced seed stream -- no published study builds a search "
    "ensemble by subsampling a common pool, and claiming so would misdescribe every method the "
    "comparison represents. Only the two hazard-filling designs need a pool, and each owns one."))
add_title(s, "The experiment at a glance")
r1y = 1.7; bw = 3.75; bh = 1.0
box(s, 0.65, r1y, bw, bh, "Stationary population\n(Kirsch–Nowak at θ_hist)", size=18)
arrow(s, 0.65 + bw, r1y + bh / 2, 0.65 + bw + 0.4, r1y + bh / 2)
box(s, 0.65 + bw + 0.4, r1y, bw, bh, "DU-forced population\n(θ ~ CMIP6 harmonic hypercube)", size=18)
arrow(s, 0.65 + 2 * bw + 0.4, r1y + bh / 2, 0.65 + 2 * bw + 0.8, r1y + bh / 2)
box(s, 0.65 + 2 * bw + 0.8, r1y, bw, bh, "Hazard metrics + redundancy\nscreen → 3–4 hazard axes", size=18)
r2y = r1y + bh + 0.45
arrow(s, 0.65 + bw + 0.4 + bw / 2, r1y + bh, 0.65 + bw * 0.75, r2y)
box(s, 0.65, r2y, bw, bh, "6 scenario designs, each\ngenerating its OWN ensemble", size=18, line=ACCENT, line_w=1.75)
arrow(s, 0.65 + bw, r2y + bh / 2, 0.65 + bw + 0.4, r2y + bh / 2)
box(s, 0.65 + bw + 0.4, r2y, bw, bh, "Independent MM-Borg searches\nN = 100 × L = 10 yr, equal NFE", size=18)
arrow(s, 0.65 + 2 * bw + 0.4, r2y + bh / 2, 0.65 + 2 * bw + 0.8, r2y + bh / 2)
box(s, 0.65 + 2 * bw + 0.8, r2y, bw, bh, "Common held-out re-evaluation\n→ multivariate satisficing", size=18)
add_body(s, [
    ("b", "**No design subsamples a shared pool** — each generates its own realizations from its own namespaced seed stream, so no two designs ever share a realization. Only the hazard-filling designs need a pool, and each owns its own."),
    ("b", "Cross-design comparison happens **once**, on the held-out test ensemble — search-time objective values are never compared across designs."),
    ("b", "Replication separates ensemble-draw variance from MOEA-seed variance; a draw re-rolls everything, including the pool."),
], top=4.75, size=19, space_after=8)

# =====================================================================
# SLIDE 15 — Testbed & optimization setup
# =====================================================================
s = add_slide(notes=(
    "Concrete machinery. FFMP parameterization = 24 decision variables; ffmp_N variants add storage "
    "zones for RQ3. Search uses a trimmed Pywr-DRB (non-NYC reservoir releases pre-simulated) for "
    "speed; the full model is used for baselines and re-evaluation. MM-Borg on Anvil; a scaling "
    "experiment currently running sets the production geometry and SU budget."))
add_title(s, "Optimization setup")
add_body(s, [
    ("b", "**Policy:** FFMP rule structure, 24 decision variables (zone-based release schedules + drought step-downs); variable-resolution variants ffmp_N add storage zones (RQ3)."),
    ("b", "**Simulation:** trimmed Pywr-DRB (non-NYC reservoir releases pre-simulated) during search; full model for baselines and re-evaluation."),
    ("b", "**Optimizer:** multi-master Borg MOEA — auto-adaptive operators, ε-dominance archive; islands × workers via MPI on Purdue Anvil."),
    ("b", "**7 objectives** spanning the basin's stakeholder axes (next slide); ε-dominance keeps 7-objective search tractable."),
    ("b", "Scaling experiment (running) sets the campaign's parallel geometry and compute budget."),
    ("gap", 2),
    ("c", "Hadka & Reed 2013, 2015; Reed et al. 2013; Hamilton et al. 2024."),
], top=1.6, size=20)

# =====================================================================
# SLIDE 16 — Objectives
# =====================================================================
s = add_slide(notes=(
    "One slide for all seven objectives; full formulations and rationale per objective live in the "
    "supplement. Two framing points: goalposts are the STATIC 1954 Decree quantities, never the live "
    "FFMP drought step-down targets (a policy must not lower its own goalpost); and noisy worst-case "
    "extremes were replaced with stable tail forms (CVaR / percentiles / counts) per Quinn et al. "
    "2017."))
add_title(s, "Objectives — seven, spanning basin stakeholders")
rows = [
    ["", "Objective", "Temporal metric (single trace)", "Dir"],
    ["1", "NYC delivery reliability", "fraction of weeks delivery ≥ 99% of Decree-capped demand", "max"],
    ["2", "NYC delivery deficit", "CVaR₉₀ of weekly deficit (% of 800 MGD)", "min"],
    ["3", "Montague flow reliability", "fraction of weeks ≥ 1,750 cfs Decree target", "max"],
    ["4", "Montague flow deficit", "CVaR₉₀ of weekly deficit (% of target)", "min"],
    ["5", "Trenton flow reliability", "fraction of weeks ≥ Decree target", "max"],
    ["6", "Downstream flood days", "days any tail gauge ≥ NWS minor flood stage", "min"],
    ["7", "NYC storage resilience", "5th percentile of daily storage (% capacity)", "max"],
]
add_table(s, rows, [0.5, 3.5, 7.15, 0.95], left=0.6, top=1.5, row_h=0.46, size=18)
add_body(s, [
    ("b", "Goalposts are the **static 1954 Decree quantities** — never the live FFMP step-down targets, which a policy could lower for itself."),
    ("b", "Stable tail / count forms replace noisy worst-case extremes; Trenton flow stands in for salt-front position (physically redundant, cleaner signal)."),
    ("c", "Hashimoto et al. 1982; Quinn et al. 2017; Rockafellar & Uryasev 2000; full definitions in supplement."),
], top=5.35, size=19, space_after=6)

# =====================================================================
# SLIDE 17 — Ensemble aggregation
# =====================================================================
s = add_slide(notes=(
    "How a (realizations x time) simulation becomes one scalar per objective during ensemble search. "
    "Two-layer annual-unit scheme, decided July 2026: stage 1 computes each objective's annual metric "
    "on every (realization x water-year) unit; stage 2 applies a per-objective operator across the "
    "pooled unit-years. Every operator follows published search-time practice (WaterPaths lineage; "
    "Quinn's 1000 annual units). Same scheme for every ensemble design -> identical unit denominator; "
    "the historic design keeps its single-trace metrics. Failure criteria / epsilons come from the "
    "sensitivity experiment."))
add_title(s, "Scoring a policy on an ensemble — two-layer aggregation")
dy = 1.6; bh = 0.95
box(s, 0.65, dy, 3.1, bh, "simulate N realizations\n× L water-years", size=18)
arrow(s, 3.75, dy + bh / 2, 4.3, dy + bh / 2)
box(s, 4.3, dy, 3.35, bh, "stage 1: annual metric per\n(realization × year) unit", size=18)
arrow(s, 7.65, dy + bh / 2, 8.2, dy + bh / 2)
box(s, 8.2, dy, 3.35, bh, "stage 2: unit operator over\npooled N·L unit-years", size=18)
arrow(s, 11.55, dy + bh / 2, 12.0, dy + bh / 2)
box(s, 12.0, dy, 0.75, bh, "fᵢ", size=20, font=TITLE_FONT)
add_body(s, [
    ("b", "**Reliability objectives** → frequency of non-failure years (a year fails if any week misses its Decree criterion) — the satisficing form used in search throughout the WaterPaths lineage."),
    ("b", "**Deficit / storage objectives** → worst 1st-percentile unit-year — stable tail emphasis over ~10³ pooled units."),
    ("b", "**Flood days** → mean annual count (P99 variant registered; expectation can mask floods)."),
    ("b", "One scheme for **every** ensemble design → identical unit denominator; the historic trace keeps its single-record metrics. Failure criteria & ε values set by a sensitivity experiment (in progress)."),
    ("gap", 2),
    ("c", "Zeff et al. 2014; Trindade et al. 2017; Quinn et al. 2017, 2018; Hamilton et al. 2022; Gold et al. 2023."),
], top=2.9, size=20)

# =====================================================================
# SLIDE 18 — Scenario designs compared
# =====================================================================
s = add_slide(notes=(
    "The heart of RQ1: six designs across TWO populations, each built by its OWN published recipe "
    "from its OWN seed stream. Hazard-filling runs in both populations, which is what makes every "
    "contrast exactly controlled: fixed_probabilistic -> hazard_filling_stationary varies only the "
    "selection RULE; input_stratified -> hazard_filling_du varies only the selection SPACE (theta vs "
    "hazard) and is the central claim. A stationary-only pool would leave input_stratified with no "
    "input space to stratify; a DU-only pool would leave hazard-filling with no exact "
    "random-selection control. resampled_probabilistic: Brodeur 2020 is the primary anchor; "
    "Trindade/Gold are cited only for the PRINCIPLE of per-evaluation re-randomization (their theta "
    "is paired with the flows and re-shuffled; ours is fused in at generation, so there is no "
    "pairing to re-randomize). Historic can't be size-matched; it is a prevailing-practice "
    "reference."))
add_title(s, "Scenario designs compared (RQ1)")
rows = [
    ["Population", "Design", "Construction of the search ensemble", "Precedent"],
    ["stationary", "historic", "the observed record as one continuous trace (N = 1)", "Giuliani et al. 2016; Herman et al. 2020"],
    ["stationary", "fixed_probabilistic", "N × L realizations generated i.i.d.; frozen across the search", "Quinn et al. 2017; Zatarain Salazar et al. 2017"],
    ["stationary", "resampled_probabilistic", "its own pool; N redrawn at every function evaluation", "Brodeur et al. 2020 (Trindade 2017 / Gold 2023: principle only)"],
    ["stationary", "**hazard_filling_stationary**", "LHS + nearest-neighbour snap in hazard space, over its own stationary pool", "**proposed**"],
    ["DU-forced", "input_stratified", "LHS over the harmonic forcing parameters; R realizations generated per design point", "Quinn et al. 2020; Bartholomew & Kwakkel 2020"],
    ["DU-forced", "**hazard_filling_du**", "the same selector, over its own DU-forced pool", "**proposed**"],
]
add_table(s, rows, [1.35, 2.85, 4.75, 3.2], left=0.6, top=1.45, row_h=0.66, size=18)
add_body(s, [
    ("b", "Every design generates its own realizations from its own seed stream; **no design is subsampled from a shared pool**."),
    ("b", "Running hazard-filling in **both** populations gives every contrast an exact within-population control; all matched designs run at **N = 100, L = 10 yr**."),
], top=6.05, size=18, space_after=2)

# =====================================================================
# SLIDE 19 — Candidate pools & forcing space
# =====================================================================
s = add_slide(notes=(
    "How the two candidate pools are built. Forcing space: CMIP6-anchored hypercube over "
    "fixed-phase harmonic amplitudes of monthly change factors (annual volume, seasonal amplitude, "
    "shoulder shape) - Quinn et al. 2018 lineage, bounded by the empirical 90% CMIP6 envelope. "
    "CRITICAL: the pools are sampled i.i.d., NEVER by LHS. A uniform random size-N subset of an "
    "i.i.d. pool has exactly the law of N fresh i.i.d. draws, which is what makes "
    "fixed_probabilistic the EXACT control for hazard_filling_stationary; a random subset of an LHS "
    "design is not i.i.d., so an LHS pool would silently void the control. input_stratified is the "
    "only design that uses LHS, and it uses it to GENERATE. Storage trick: keep only the hazard "
    "image + seeds; realizations regenerate deterministically (~440 GB avoided). Scope caveat: "
    "historical interannual persistence is retained; claims scoped accordingly."))
add_title(s, "Candidate pools & the deeply uncertain forcing space")
add_body(s, [
    ("b", "**Forcing space:** CMIP6-anchored hypercube over harmonic amplitudes of monthly streamflow change factors — annual volume, seasonal amplitude, shoulder shape — bounded by the empirical 90% CMIP6 envelope (fixed seasonal phases)."),
    ("b", "**Two pools, one per hazard-filling design:** stationary (Kirsch–Nowak at the historic fit) and DU-forced (θ from the hypercube) → ~10⁵–10⁶ sequences each. Each design owns its pool; the pool is re-drawn on every ensemble draw."),
    ("b", "**Pools are sampled i.i.d., never by LHS** — a random subset of an LHS design is not i.i.d., which would silently void the exact control. Only input_stratified uses LHS, and it uses it to *generate*, never to build a pool."),
    ("b", "**Windows:** disjoint 10-yr scenarios — long enough to contain the 1960s drought of record plus onset and recovery; fixed initial storage with a 1-yr warm-up."),
    ("b", "**Storage:** only hazard coordinates + seeds persist; any realization regenerates deterministically on demand (verified) — avoids ~440 GB of traces. Seed domains are disjoint across designs, draws, and the test ensemble."),
    ("b", "**Scope:** historical interannual persistence retained; claims scoped accordingly."),
    ("gap", 2),
    ("c", "Kirsch et al. 2013; Nowak et al. 2010; Quinn et al. 2018."),
], top=1.6, size=19, space_after=8)

# =====================================================================
# SLIDE 20 — Hazard axes & selector
# =====================================================================
s = add_slide(notes=(
    "Hazard axes: SSI-6 drought event metrics (run theory) plus peaks-over-threshold flood metrics. "
    "A redundancy screen (Spearman clustering at |rho|>=0.7, tail-balanced) selects 3-4 "
    "low-collinearity axes on each production pool's hazard image. Selector: normalize each axis by "
    "its empirical CDF, draw N Latin-hypercube anchors, snap each to the nearest not-yet-selected "
    "pool member. Deterministic given the anchor seed. NO annealing, no tuning, no discrepancy "
    "objective -- which is exactly why L2-star discrepancy stays an INDEPENDENT diagnostic of the "
    "achieved design rather than the quantity being minimized, so coverage is reported as a RESULT "
    "and not merely as a build-verification gate. Be upfront: this deliberately distorts scenario "
    "probabilities; that is the design choice under study, and the held-out re-evaluation is the "
    "corrective."))
add_title(s, "Hazard axes & the hazard-filling selector")
add_body(s, [
    ("b", "**Candidate axes:** SSI-6 drought-event metrics (deficit volume, duration, peak depth, onset and recovery rate) and peaks-over-threshold flood metrics (peak magnitude, pulse duration, rise rate) — computed per sequence."),
    ("b", "**Redundancy screen** on each pool's hazard image — Spearman clustering (|ρ| ≥ 0.7 = redundant), tail-balanced across the dry and wet concepts → **3–4** low-collinearity axes."),
    ("b", "**Selector:** normalize each axis by its empirical CDF, draw **N Latin-hypercube anchors**, snap each to the **nearest unused pool member**. Deterministic given the anchor seed; **no annealing, no tuning**. The snap is intrinsic — hazard coordinates cannot be prescribed at generation."),
    ("b", "Because the selector does **not** optimize a discrepancy objective, **L2-star discrepancy stays an independent diagnostic** of the achieved design — coverage is a **result**, reported against a random design at the same (N, m), not a build-verification gate."),
    ("b", "**Deliberate probability distortion:** severe corners are over-represented relative to their frequency — this is the design choice under study; the held-out re-evaluation is the probability-faithful corrective."),
    ("gap", 2),
    ("c", "Yevjevich 1967; Vicente-Serrano et al. 2012; Richter et al. 1996; Olden & Poff 2003; Fang et al. 2000; Bonham et al. 2024."),
], top=1.6, size=19, space_after=8)

# =====================================================================
# SLIDE 21 — Controls for a fair comparison
# =====================================================================
s = add_slide(notes=(
    "The credibility slide. The budget control is a CONSEQUENCE of the sizing choice, not something "
    "imposed: because N = 100 and L = 10 yr are common to every matched design, per-evaluation "
    "cost, warm-up, scenario-years and wall-clock are identical, so equal-NFE and equal-"
    "scenario-years coincide. ONE budget condition, no arms, no composition-vs-search-effort "
    "confound. A common (N, L) is REQUIRED, not convenient: with different L the selection rule "
    "would be confounded with record length. Replication: a draw re-rolls EVERYTHING random about "
    "building the ensemble, including the pool -- otherwise a hazard-filling draw would vary only "
    "its anchor plan while a fixed_probabilistic draw re-rolls its whole sample, and hazard-filling "
    "would look more stable by construction rather than as a finding."))
add_title(s, "Controls for a fair comparison")
add_body(s, [
    ("b", "**Budget:** every matched design runs at **N = 100, L = 10 yr** → 1,000 scenario-years per evaluation, at **equal NFE**. Because N and L are common, per-evaluation cost, warm-up, scenario-years and wall-clock are *identical* — equal-NFE and equal-scenario-years **coincide**. One budget condition, no arms, no confound between composition and search effort."),
    ("b", "The common (N, L) is **required, not convenient**: if L differed across designs, the selection rule would be confounded with record length. N is bounded below by the fill requirement (~3.2 points per dimension at m = 4)."),
    ("b", "**Replication:** K ≈ 10 independent ensemble draws × S = 2–3 MOEA seeds per design; a **draw** is the design's construction re-run from scratch with a fresh seed, **including its pool**. The draw is the **unit of analysis** (effective n ≈ K, not K·S)."),
    ("b", "**Seed-stream disjointness:** every design, every draw, and the test ensemble generate from a namespaced seed domain, so no two ever share a realization."),
    ("b", "Search-time objective values are **never compared across designs** — they are computed on different ensembles; convergence diagnostics are reported per design only."),
    ("gap", 2),
    ("c", "Zatarain Salazar et al. 2017 (size–fidelity–cost); Kaut & Wallace 2007 (stability criteria)."),
], top=1.6, size=18, space_after=8)

# =====================================================================
# SLIDE 22 — Evaluation on the held-out test ensemble
# =====================================================================
s = add_slide(notes=(
    "The single point of cross-design comparison, and the finalized metric set (src/robustness.py). "
    "PRIMARY: the multivariate Starr (1962) domain criterion -- realization-level, the standard "
    "unit of the Herman/Trindade/Gold lineage, and it converges at 50-300 scenarios (Bonham 2024). "
    "Everything else is free from the same persisted (solution x realization x objective) cube. "
    "Regret is NOT computed, and no perfect-foresight optimization appears anywhere in the study: "
    "Cohen (2021) baseline regret needs one perfect-foresight MOEA run per scenario (97 "
    "optimizations, ~3,233 CPU-h), and regret-from-best is set-relative and design-coupled -- "
    "dropping one design changes every other design's score -- and never converges on a tail "
    "objective. Pooled-reference-set hypervolume is likewise not a cross-design yardstick "
    "(Zatarain Salazar 2017 sec 5.3)."))
add_title(s, "Evaluation — one common held-out test ensemble")
add_body(s, [
    ("b", "**Test ensemble:** ≥1,000 long (≥30-yr) realizations spanning the full forcing space; never used during any search and **never the source of any search ensemble**; ≥2 constructions (one from a structurally different generator) bound how rankings depend on its design."),
    ("b", "Every design's final Pareto-approximate set is re-simulated on this ensemble; nondominated sets are recomputed from re-evaluated values for all designs alike."),
    ("b", "**Primary metric: multivariate satisficing** (Starr 1962 domain criterion) — the fraction of test realizations meeting **all** criteria jointly; the realization is the unit."),
    ("b", "**Secondary, all free from the same cube:** univariate satisficing; Laplace/mean and maximin (risk-neutral and risk-averse anchors); improvement over the status quo (vs the current FFMP policy — a fixed external reference, no optimization); an attainability screen for realizations no policy can win."),
    ("b", "**Not used: regret, and no perfect-foresight optimization anywhere** — regret-from-best is set-relative and design-coupled, and does not converge on a tail objective. **Not used: pooled-reference-set hypervolume** as a cross-design measure."),
    ("b", "Ranking agreement via Kendall's τ_b **across the design rankings**; raw performance distributions co-reported with every robustness number (the degeneracy check)."),
    ("gap", 2),
    ("c", "Starr 1962; Herman et al. 2015; McPhail et al. 2018, 2020; Kasprzyk et al. 2013; Shavazipour et al. 2021; Bonham et al. 2024; Zatarain Salazar et al. 2017."),
], top=1.6, size=18, space_after=7)

# =====================================================================
# SLIDE 23 — Status & next steps
# =====================================================================
s = add_slide(notes=(
    "Where things stand as of July 2026. Built and verified: the full pipeline runs end-to-end at "
    "test scale; per-design generation with disjoint seed domains; deterministic pools with "
    "single-realization regeneration; the LHS + nearest-neighbour hazard-filling selector is wired; "
    "the Anvil scaling experiment is running now and gates the campaign scope. The sensitivity "
    "experiments finalize thresholds/epsilons before the campaign."))
add_title(s, "Status & next steps")
add_body(s, [
    ("h", "Built and verified"),
    ("b", "End-to-end pipeline (generation → staging → MM-Borg search → re-evaluation) runs at test scale; historic-design production runs launchable on Anvil."),
    ("b", "Per-design generation with disjoint seed domains; deterministic candidate pools with on-demand single-realization regeneration; hazard-filling selector (LHS + nearest-neighbour) wired; Anvil scaling experiment currently running."),
    ("gap", 6),
    ("h", "Next"),
    ("b", "Generate the production candidate pools (stationary and DU-forced); run the hazard-axis redundancy screen."),
    ("b", "Objective-sensitivity experiments → finalize failure criteria, flood operator, ε values."),
    ("b", "Stage the held-out test ensemble; run the HPC campaign (designs × draws × seeds); re-evaluate and analyze."),
], top=1.6, size=20)

# =====================================================================
# SLIDE 24 — Feedback requested
# =====================================================================
s = add_slide(notes=(
    "Close by asking for input on the genuinely open decisions, in priority order. (1) Campaign "
    "scope vs compute budget - the SU table lands when the scaling experiment finishes. (2) Test-"
    "ensemble design, including whether one construction should use a structurally different "
    "generator. (3) The hazard-axis set after the redundancy screen. (4) Flood unit operator and "
    "annual failure criteria from the sensitivity experiment. (5) The centre and span of the "
    "satisficing-threshold grid, which the main-text threshold sweep is run over."))
add_title(s, "Open decisions — committee input requested")
add_body(s, [
    ("b", "**Campaign scope vs compute budget** — how many designs × draws × seeds to run (SU costs pending the Anvil scaling results)."),
    ("b", "**Test-ensemble design** — breadth of the deeply uncertain envelope; which structurally different streamflow generator for the second construction?"),
    ("b", "**Hazard-axis set** — which 3–4 axes after the redundancy screen on the production candidate pools; the N_θ / R split for input_stratified."),
    ("b", "**Objective details** — flood unit operator (mean vs P99) and annual failure criteria, from the objective-sensitivity experiment."),
    ("b", "**The satisficing-threshold grid** — its centre and span, fixed before any cross-design comparison."),
    ("gap", 8),
    ("p", "Thank you — discussion welcome."),
], top=1.7, size=20)

# =====================================================================
# SUPPLEMENTAL
# =====================================================================
s = add_slide(notes="Supplemental divider.")
divider(s, [("Supplemental", True)])

# ---- S: Notation ----
s = add_slide(notes=(
    "Plain-language building blocks used by every objective slide that follows."))
add_title(s, "Notation — plain-language building blocks")
add_body(s, [
    ("f", "shortfall:  (gap)₊ = max(gap, 0) — the part of a miss below the goal; never negative"),
    ("f", "weekly value = the average (or total) of that week's daily values"),
    ("f", "reliability = (weeks that meet the goal) ÷ (total weeks)"),
    ("f", "CVaR₉₀(series) = the average of the worst 10% of the weekly values"),
    ("f", "P₅(series) = the 5th-percentile value (only 5% of days fall below it)"),
], top=1.5, size=20, space_after=8)
add_body(s, [
    ("b", "Flows and demands in MGD (million gallons per day); storage in MG."),
    ("b", "Decree goalposts (fixed): NYC diversion cap 800 MGD; Montague 1,131.05 MGD (= 1,750 cfs); Trenton 1,938.95 MGD. NYC system capacity = 270,837 MG (Cannonsville + Pepacton + Neversink)."),
    ("b", "Objectives are computed on the daily series after dropping the first 365 days (model warm-up)."),
], top=4.35, size=19)

# ---- Objective detail slides ----
def objective_slide(num, name, direction, formulas, bullets, cites, notes):
    s = add_slide(notes=notes)
    add_title(s, "Objective %d — %s  (%s)" % (num, name, direction))
    add_body(s, [("f", f) for f in formulas], top=1.5, size=20, space_after=6)
    items = [("b", b) for b in bullets]
    items += [("gap", 2), ("c", cites)]
    add_body(s, items, top=1.5 + 0.5 * len(formulas) + 0.25, size=19)
    return s

objective_slide(
    1, "NYC delivery reliability", "higher is better",
    ["demand = min(NYC demand, 800 MGD)      [capped at the Decree right]",
     "f₁ = (weeks with delivery ≥ 99% of demand) ÷ (total weeks)"],
    ["How often NYC receives essentially all the water it is entitled to draw — a weekly success rate comparing weekly totals of delivery and demand.",
     "Demand is capped at the 800 MGD Decree right so voluntary winter low-takes are not counted as shortfalls; only forced failures count against the policy.",
     "A success rate is a stable, fast-converging signal for the optimizer, unlike noisier magnitude statistics.",
     "Ensemble search form: frequency of non-failure years (a year fails if any week misses the criterion)."],
    "Hashimoto et al. 1982; Herman et al. 2015; Zeff et al. 2014; Bonham et al. 2024.",
    "Weekly satisficing frequency; Decree-capped demand; annual failure-year frequency in ensemble mode.")

objective_slide(
    2, "NYC delivery deficit", "lower is better",
    ["weekly deficit = max(demand − delivery, 0) ÷ 800 MGD × 100%",
     "f₂ = average of the worst 10% of weekly deficits      (CVaR₉₀)"],
    ["How bad shortfalls get, not just how often: the average size of the worst 10% of weekly supply gaps, as a percent of the 800 MGD cap.",
     "Scaling to the fixed cap makes a 50 MGD shortfall read the same in any season.",
     "Worst-10% average instead of the single worst week: one extreme week is noisy and irreproducible; averaging the tail keeps the focus on bad outcomes while staying stable.",
     "Ensemble search form: worst 1st-percentile unit-year of the within-year CVaR₉₀."],
    "Rockafellar & Uryasev 2000; Quinn et al. 2017; Fairbrother et al. 2022.",
    "CVaR90 tail deficit; ensemble mode uses worst-P99 unit-year.")

objective_slide(
    3, "Montague flow reliability", "higher is better",
    ["f₃ = (weeks with Montague flow ≥ 1,131.05 MGD) ÷ (total weeks)",
     "1,131.05 MGD = 1,750 cfs      (1954 Decree target)"],
    ["How often flow at Montague meets NYC's downstream obligation — a weekly success rate on weekly-average flow.",
     "Will not reach 100%: FFMP drought rules deliberately cut releases below target in drought, so the achievable ceiling depends on the policy — the signal stays continuous.",
     "Scored against the fixed Decree target, never the live FFMP target, to avoid rewarding self-lowering drought step-downs.",
     "Ensemble search form: frequency of non-failure years."],
    "Hashimoto et al. 1982; Herman et al. 2015; Zeff et al. 2014.",
    "Weekly reliability vs static Decree target.")

objective_slide(
    4, "Montague flow deficit", "lower is better",
    ["weekly deficit = max(1,131.05 − flow, 0) ÷ 1,131.05 × 100%",
     "f₄ = average of the worst 10% of weekly deficits      (CVaR₉₀)"],
    ["How far Montague flow falls below target in the bad weeks, as a percent of the target.",
     "Montague flow is storm-dominated, so its single worst week is mostly exogenous noise — the tail average matters most here.",
     "Ensemble search form: worst 1st-percentile unit-year of the within-year CVaR₉₀."],
    "Rockafellar & Uryasev 2000; Quinn et al. 2017.",
    "CVaR90 of Montague deficit.")

objective_slide(
    5, "Trenton flow reliability", "higher is better",
    ["f₅ = (weeks with Trenton flow ≥ 1,938.95 MGD) ÷ (total weeks)"],
    ["How often flow at Trenton meets the lower-basin / New Jersey obligation — a weekly success rate against the fixed target.",
     "Replaces an explicit salt-front objective: the Trenton target exists largely to repel salt intrusion (physically redundant), and Trenton flow is a clean signal where the salt-front model is unreliable in extreme drought.",
     "Gives New Jersey, a co-equal Decree party, direct representation so the search can discover NYC ↔ NJ conflicts.",
     "Ensemble search form: frequency of non-failure years."],
    "Hashimoto et al. 1982; Herman et al. 2015; Trindade et al. 2017; Hadjimichael et al. 2020.",
    "Trenton reliability replacing salt-front objective.")

objective_slide(
    6, "Downstream flood days", "lower is better",
    ["f₆ = days on which any tail gauge reaches its NWS minor flood stage",
     "tail gauges:  Hale Eddy, Fishs Eddy, Bridgeville"],
    ["A simple count of flood-onset days just below the reservoirs: a day counts if any of the three tail gauges hits its minor (flood-onset) stage.",
     "These gauges respond to NYC's release decisions — unlike storm-dominated mainstem flow — so flooding here is attributable to operations.",
     "Counting threshold exceedances avoids the expected-damage trap of multiplying probabilities by an assumed damage curve.",
     "Ensemble search form: mean annual flood days (P99 variant registered — expectation can mask floods)."],
    "Quinn et al. 2017.",
    "Flood day count at minor stage; ensemble operator pending sensitivity experiment.")

objective_slide(
    7, "NYC storage resilience", "higher is better",
    ["storage% = combined NYC storage ÷ 270,837 MG × 100%",
     "f₇ = the 5th percentile of daily storage%      (P₅)"],
    ["How low NYC storage routinely gets: the level combined storage stays above 95% of the time, as a percent of full capacity.",
     "The single-day minimum is dominated by one drought event and is noisy; a low percentile is reproducible across scenarios.",
     "Ensemble search form: 1st-percentile unit-year of the annual minimum storage."],
    "Quinn et al. 2017.",
    "Storage P5; ensemble mode uses P01 of annual minima.")

# ---- S: taxonomy ----
s = add_slide(notes=(
    "Condensed version of the four-family taxonomy (working note). The proposed design sits in "
    "family III (bottom-up coverage), sub-family 'coverage-based subset selection from realized "
    "sequences' - distinguished by scale (a 10^5-10^6 candidate pool), multi-dimensional hazard "
    "coordinates, and destination (the search ensemble). Family IV is where Cohen (2021) sits: it "
    "selects by decision relevance, which costs an optimization per scenario -- we cite that "
    "finding as motivation and do not compute any such quantity."))
add_title(s, "Supplemental — taxonomy of scenario-design families")
rows = [
    ["Family (criterion)", "Examples", "Design space"],
    ["I. Judgment-based (expert precedent)", "historical trace; design droughts; storylines", "single traces / events"],
    ["II. Distribution-driven (fidelity to a distribution)", "SAA random sampling; LHS on inputs; scenario reduction; distribution-matching subset selection", "input / probability space"],
    ["III. Bottom-up coverage (span a condition space)", "exposure-space grids; SOW ensembles; controlled generation; **coverage-based subset selection ← this study**", "forcing attributes → realized flow metrics"],
    ["IV. Problem-driven (relevance to the decision)", "importance sampling; decision-relevance selection (Cohen 2021 — cited as motivation, never computed here); adaptive selection", "outcome-informed"],
]
add_table(s, rows, [4.5, 5.35, 2.3], left=0.6, top=1.55, row_h=0.62, size=18)
add_body(s, [
    ("c", "Brown et al. 2012; Fairbrother et al. 2022; Fowler et al. 2024; full taxonomy in working notes."),
], top=6.0, size=18)

# ---- S: forcing parameterization ----
s = add_slide(notes=(
    "Detail for the forcing hypercube. Each CMIP6 (GCM x SSP x period) run gives a 12-month "
    "multiplicative change-factor profile for NYC inflows; its log profile is decomposed into a "
    "low-order harmonic with phases FIXED at the canonical CMIP6 seasonal shape; only amplitudes "
    "(annual mean m, annual r1, semiannual r2) are sampled, by LHS, within the empirical 90% "
    "envelope. Factors adjust the Kirsch generator's log-space monthly mean/SD via Kirsch et al. "
    "2013 eqs. 10-11. An optional CV axis perturbs variance independently."))
add_title(s, "Supplemental — forcing-space parameterization")
add_body(s, [
    ("b", "Each CMIP6 run (GCM × SSP × period) → a 12-month multiplicative change-factor profile for NYC inflow gages."),
    ("b", "Log-profile decomposed into a low-order harmonic; **phases fixed at the canonical CMIP6 seasonal shape** (correct winter peak + asymmetry), amplitudes sampled."),
    ("b", "Sampled axes: annual-mean level m, annual amplitude r₁, semiannual amplitude r₂ (volume / seasonal amplitude / shoulder shape) + optional variance axis."),
    ("b", "LHS within the empirical 90% CMIP6 envelope → deeply uncertain magnitude combinations the GCMs did not jointly produce, with a physically plausible seasonal shape."),
    ("b", "Factors applied to the Kirsch generator's log-space monthly mean / SD (Kirsch et al. 2013, eqs. 10–11) before generation."),
    ("gap", 2),
    ("c", "Quinn et al. 2018 (harmonic lineage); Kirsch et al. 2013."),
], top=1.6, size=20)

# ---- S: Anvil scaling ----
s = add_slide(notes=(
    "Supplemental: the scaling experiment currently running on Anvil that sets the production "
    "geometry and SU budget. Two parts: a packing sweep over ensemble-evaluation parallelism "
    "(realizations per Pywr-DRB simulation) and a Borg strong-scaling ladder (islands x workers)."))
add_title(s, "Supplemental — Anvil scaling experiment (running)")
add_body(s, [
    ("b", "Purpose: fix the production parallel geometry and the campaign's SU budget before committing the RQ1 campaign."),
    ("h", "Parallelism levers"),
    ("b", "Nodes and cores per node; number of MM-Borg islands; workers per island; realizations evaluated in parallel within one Pywr-DRB simulation."),
    ("h", "Design"),
    ("b", "Packing sweep: throughput vs realizations-per-simulation (ensemble evaluation path)."),
    ("b", "Borg strong scaling: fixed problem, island × worker ladder at short NFE."),
    ("b", "Output: SU cost per (design × draw × seed) → campaign scope table."),
], top=1.6, size=20)

# ---------- save ----------
OUT = r"c:\Users\tjame\Desktop\Research\DRB\Pywr-DRB\NYCOptimization\docs\NYCOpt_Objective_Definitions_CLAUDE.pptx"
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
